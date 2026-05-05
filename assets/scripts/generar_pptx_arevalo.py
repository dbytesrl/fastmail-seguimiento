"""
Generador de presentación PowerPoint — Arevalo Tucumán × Boreal Marketing
Audit de presencia digital + bases para rediseño y campañas
Ejecutar: python3 generar_pptx_arevalo.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PALETA BOREAL ────────────────────────────────────────────────
VIOLETA     = RGBColor(0x33, 0x00, 0xCC)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO       = RGBColor(0x00, 0x00, 0x00)
CELESTE     = RGBColor(0x12, 0xBD, 0xFF)
ROSA        = RGBColor(0xFE, 0x67, 0xB4)
GRIS_CLARO  = RGBColor(0xF5, 0xF5, 0xF5)
GRIS_MED    = RGBColor(0xA0, 0xA0, 0xA0)
GRIS_TEXTO  = RGBColor(0x33, 0x33, 0x33)
VIOLETA_OSC = RGBColor(0x22, 0x00, 0x99)
VIOLETA_MED = RGBColor(0x2A, 0x00, 0xBB)
VERDE_OK    = RGBColor(0x00, 0xB0, 0x60)
ROJO_NOK    = RGBColor(0xCC, 0x22, 0x22)
NARANJA     = RGBColor(0xFF, 0x88, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════

def rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h,
        size=14, bold=False, italic=False,
        color=BLANCO, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Montserrat"
    return tb

def footer(slide):
    rect(slide, 0, H - Inches(0.36), W, Inches(0.36), VIOLETA)
    txt(slide, "Confidencial · Boreal Marketing × Arevalo Tucumán · Auditoría Digital 2026",
        Inches(0.3), H - Inches(0.34), Inches(9), Inches(0.28),
        size=8, color=BLANCO)
    txt(slide, "boreal.com.ar",
        W - Inches(2), H - Inches(0.34), Inches(1.8), Inches(0.28),
        size=8, color=CELESTE, align=PP_ALIGN.RIGHT)

def header(slide, titulo, sub=None):
    rect(slide, 0, 0, W, Inches(1.05), VIOLETA)
    txt(slide, titulo,
        Inches(0.5), Inches(0.15), W - Inches(1), Inches(0.6),
        size=24, bold=True, color=BLANCO)
    if sub:
        txt(slide, sub,
            Inches(0.5), Inches(0.7), W - Inches(1), Inches(0.3),
            size=11, italic=True, color=CELESTE)

def pill(slide, label, l, t, color=VIOLETA, text_color=BLANCO, size=10):
    rect(slide, l, t, Inches(1.5), Inches(0.3), color)
    txt(slide, label, l + Inches(0.08), t + Inches(0.04),
        Inches(1.35), Inches(0.24),
        size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def tag_critico(slide, l, t):
    pill(slide, "🔴 CRÍTICO", l, t, ROJO_NOK)

def tag_alto(slide, l, t):
    pill(slide, "🟠 ALTO", l, t, NARANJA)

def tag_medio(slide, l, t):
    pill(slide, "🟡 MEDIO", l, t, RGBColor(0xCC, 0xAA, 0x00))

def divider(slide, y, color=CELESTE):
    ln = slide.shapes.add_shape(1, Inches(0.5), y, W - Inches(1), Pt(2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def row_tabla(slide, cols_data, cols_x, cols_w, y, row_h=Inches(0.5), bg=GRIS_CLARO, text_color=NEGRO, size=11):
    for val, x, w in zip(cols_data, cols_x, cols_w):
        rect(slide, x, y, w - Inches(0.04), row_h, bg)
        txt(slide, str(val), x + Inches(0.1), y + Inches(0.07),
            w - Inches(0.2), row_h - Inches(0.1),
            size=size, color=text_color, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)
rect(s, W - Inches(5.5), 0, Inches(5.5), H * 0.52, ROSA)
rect(s, W - Inches(4.5), H * 0.52, Inches(4.5), H * 0.48, CELESTE)
rect(s, W - Inches(2.8), Inches(1.8), Inches(2.8), Inches(2.8), VIOLETA_OSC)

txt(s, "BOREAL MARKETING",
    Inches(0.6), Inches(0.38), Inches(5), Inches(0.42),
    size=11, bold=True, color=CELESTE)
rect(s, Inches(0.6), Inches(0.84), Inches(2.5), Pt(2.5), ROSA)

txt(s, "Auditoría Digital",
    Inches(0.6), Inches(1.0), Inches(8.5), Inches(0.65),
    size=36, bold=False, italic=True, color=BLANCO)
txt(s, "Arevalo Tucumán",
    Inches(0.6), Inches(1.6), Inches(9), Inches(1.0),
    size=52, bold=True, color=BLANCO)

txt(s, "Análisis de diseño · Estructura · Contenido\nBases para rediseño y campañas de Ads + Redes",
    Inches(0.6), Inches(2.75), Inches(8.2), Inches(0.9),
    size=15, italic=True, color=BLANCO)

rect(s, Inches(0.6), Inches(3.85), Inches(7), Pt(2), ROSA)

txt(s, "Preparado por Boreal Marketing · Marzo 2026",
    Inches(0.6), Inches(3.95), Inches(7), Inches(0.38),
    size=11, color=GRIS_MED)
txt(s, "arevalotucuman.ar",
    Inches(0.6), Inches(4.38), Inches(4), Inches(0.35),
    size=11, color=CELESTE)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — RESUMEN EJECUTIVO
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "RESUMEN EJECUTIVO",
       "Lo que encontramos en una mirada")

rect(s, Inches(0.4), Inches(1.18), W - Inches(0.8), Inches(0.88), VIOLETA)
txt(s, "Arevalo tiene los activos correctos para dominar el mercado de salud prepaga en Tucumán. "
       "El problema no es el servicio — es cómo está comunicado y organizado digitalmente.",
    Inches(0.6), Inches(1.26), W - Inches(1.2), Inches(0.72),
    size=14, bold=True, italic=True, color=CELESTE, align=PP_ALIGN.CENTER)

items = [
    (VERDE_OK,  "LO QUE FUNCIONA",
     ["5 sucursales en toda la provincia",
      "24 especialidades médicas disponibles",
      "Meta Pixel instalado y activo",
      "App propia en Google Play",
      "WhatsApp flotante con 2 canales",
      "25+ años de trayectoria real"]),
    (NARANJA, "LO QUE HAY QUE ACTIVAR",
     ["GA4 instalado pero apagado (1 cambio lo activa)",
      "Dos audiencias mezcladas sin separación",
      "Hero con carrusel: el mayor destructor de CVR",
      "Sin precios visibles en seguros ni planes",
      "6 números de teléfono distintos en el sitio",
      "Sin testimonios de afiliados"]),
    (VIOLETA, "LA OPORTUNIDAD",
     ["Competencia digital mínima en el rubro en Tucumán",
      "Búsquedas activas sin capturar en Google",
      "Audiencias de retargeting disponibles (pixel activo)",
      "Interior de provincia sin presencia digital real",
      "SEO local con potencial inmediato",
      "Content health con alto alcance orgánico"]),
]

col_x = [Inches(0.4), Inches(4.65), Inches(8.9)]
col_w = Inches(4.0)
for (col, titulo, puntos), x in zip(items, col_x):
    rect(s, x, Inches(2.22), col_w, Inches(0.42), col)
    txt(s, titulo, x + Inches(0.12), Inches(2.26),
        col_w - Inches(0.25), Inches(0.35),
        size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    for ii, p in enumerate(puntos):
        yy = Inches(2.7) + ii * Inches(0.62)
        bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
        rect(s, x, yy, col_w, Inches(0.58), bg)
        rect(s, x, yy, Inches(0.06), Inches(0.58), col)
        txt(s, p, x + Inches(0.15), yy + Inches(0.09),
            col_w - Inches(0.28), Inches(0.42),
            size=10.5, color=NEGRO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — DIAGNÓSTICO TÉCNICO
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "DIAGNÓSTICO TÉCNICO",
       "Estado del stack tecnológico y herramientas de medición")

cx = [Inches(0.4), Inches(3.5), Inches(7.2), Inches(10.2)]
cw = [Inches(3.05), Inches(3.65), Inches(2.95), Inches(2.85)]
headers_t = ["COMPONENTE", "ESTADO", "DETALLE", "ACCIÓN REQUERIDA"]
for h_, x, w in zip(headers_t, cx, cw):
    rect(s, x, Inches(1.15), w - Inches(0.04), Inches(0.42), VIOLETA)
    txt(s, h_, x + Inches(0.1), Inches(1.17), w - Inches(0.2), Inches(0.38),
        size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

rows = [
    ("CMS",                  "✅ WordPress",          "Actualizable, manejable",                     "Mantener actualizado"),
    ("Google Analytics UA",  "⚠️ Deprecated jul/2023","Universal Analytics ya no recolecta datos",   "Reemplazar por GA4 ya"),
    ("Google Analytics 4",   "⚠️ Instalado — APAGADO","Script presente, enabledGoogle: OFF",         "🔴 Activar en config — 5 minutos"),
    ("Meta Pixel",           "✅ Activo",              "Audiencias acumulándose",                     "Configurar eventos: Lead, Contact"),
    ("Google Tag Manager",   "❌ No detectado",        "Sin capa de gestión de tags",                 "Implementar — centraliza todo"),
    ("Conversiones Goog.Ads","❌ No configurado",      "Campañas sin datos de conversión",            "Implementar vía GTM"),
    ("HTTPS / SSL",          "✅ Trust Provider",      "Certificado activo",                          "OK"),
    ("URL con typo",         "❌ /comlementos/",       "URL rota en menú de navegación",              "Corregir a /complementos/"),
    ("Twitter widget",       "⚠️ Script activo",       "Carga innecesaria si cuenta inactiva",        "Evaluar eliminar"),
    ("App Google Play",      "✅ Disponible",          "com.ionicframework.arevalo820026",             "Promover como diferencial"),
]

alt = [GRIS_CLARO, BLANCO]
for ii, row in enumerate(rows):
    y = Inches(1.64) + ii * Inches(0.48)
    bg = alt[ii % 2]
    for val, x, w in zip(row, cx, cw):
        rect(s, x, y, w - Inches(0.04), Inches(0.44), bg)
        c = NEGRO
        if "🔴" in val:    c = ROJO_NOK
        elif "✅" in val:  c = VERDE_OK
        elif "❌" in val:  c = ROJO_NOK
        elif "⚠️" in val: c = NARANJA
        txt(s, val, x + Inches(0.1), y + Inches(0.06),
            w - Inches(0.2), Inches(0.34),
            size=10, color=c, bold=("🔴" in val))

rect(s, Inches(0.4), Inches(6.52), W - Inches(0.8), Inches(0.52), RGBColor(0xFF, 0xF0, 0xD0))
rect(s, Inches(0.4), Inches(6.52), Inches(0.1), Inches(0.52), NARANJA)
txt(s, "⚡  Quick win: GA4 está instalado y apagado. Activarlo toma 5 minutos y desbloquea todos los datos de tráfico desde el primer día.",
    Inches(0.6), Inches(6.6), W - Inches(1.2), Inches(0.38),
    size=11, bold=True, italic=True, color=RGBColor(0x88, 0x44, 0x00))

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — ANÁLISIS DE DISEÑO Y UX
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "ANÁLISIS DE DISEÑO Y UX",
       "Lo que el usuario ve, siente y hace — o no hace")

# Panel izquierdo — problemas
rect(s, Inches(0.4), Inches(1.18), Inches(6.1), Inches(0.42), ROJO_NOK)
txt(s, "PROBLEMAS DETECTADOS", Inches(0.5), Inches(1.22),
    Inches(5.9), Inches(0.35), size=12, bold=True, color=BLANCO)

problemas = [
    ("Hero carrusel de 6 slides",
     "Los usuarios no esperan al slide 3. El primer frame define todo y hoy no tiene propuesta de valor clara ni CTA único."),
    ("Sin bifurcación de audiencias",
     "El usuario que quiere afiliarse ve lo mismo que el que quiere sacar un turno. Ninguno encuentra lo que busca rápido."),
    ("6 números de teléfono distintos",
     "0381 485-9000 · 381 543-4242 · 430 7878 · 430 7716 · 551-0686 · 381 208-6164. El usuario no sabe cuál usar."),
    ("Competencia de CTAs",
     '"Reservá tu turno" + "Abona tu cuota" + "Descargá la app" al mismo nivel visual. Sin jerarquía, ninguno convierte bien.'),
    ("Sin sistema visual unificado",
     "Tipografías, colores y estilos de cards mezclan distintas generaciones del sitio. Genera desconfianza subconsciente."),
]
for ii, (titulo, desc) in enumerate(problemas):
    yy = Inches(1.68) + ii * Inches(0.92)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(0.4), yy, Inches(6.1), Inches(0.88), bg)
    rect(s, Inches(0.4), yy, Inches(0.08), Inches(0.88), ROJO_NOK)
    txt(s, titulo, Inches(0.58), yy + Inches(0.06),
        Inches(5.8), Inches(0.3), size=11, bold=True, color=NEGRO)
    txt(s, desc, Inches(0.58), yy + Inches(0.36),
        Inches(5.8), Inches(0.45), size=10, italic=True, color=GRIS_TEXTO)

# Panel derecho — oportunidades de diseño
rect(s, Inches(6.8), Inches(1.18), Inches(6.1), Inches(0.42), VERDE_OK)
txt(s, "OPORTUNIDADES DE DISEÑO", Inches(6.9), Inches(1.22),
    Inches(5.9), Inches(0.35), size=12, bold=True, color=BLANCO)

oportunidades = [
    ("Hero estático con 2 CTAs primarios",
     '"Quiero afiliarme" y "Ya soy afiliado" — bifurca al usuario en el segundo 1.'),
    ("Tarjeta de emergencia siempre visible",
     "El número de ambulancias debe estar fijo en header o sticky bar. En salud, la urgencia es el momento de mayor confianza."),
    ("Sección 'Por qué Arevalo' con datos",
     "25 años · 24 especialidades · 5 sucursales · App propia. Esos números generan confianza instantánea."),
    ("Testimonios de afiliados reales",
     "En salud, el boca a boca digital es el driver #1. Una sección con 6-8 historias reales cambia el CVR."),
    ("Menú simplificado mobile-first",
     "Máximo 4 ítems en navegación principal. Todo lo demás, en sub-menús o en la home según audiencia."),
]
for ii, (titulo, desc) in enumerate(oportunidades):
    yy = Inches(1.68) + ii * Inches(0.92)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(6.8), yy, Inches(6.1), Inches(0.88), bg)
    rect(s, Inches(6.8), yy, Inches(0.08), Inches(0.88), VERDE_OK)
    txt(s, titulo, Inches(6.98), yy + Inches(0.06),
        Inches(5.8), Inches(0.3), size=11, bold=True, color=NEGRO)
    txt(s, desc, Inches(6.98), yy + Inches(0.36),
        Inches(5.8), Inches(0.45), size=10, italic=True, color=GRIS_TEXTO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — EL PROBLEMA CENTRAL DE ARQUITECTURA
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

txt(s, "EL PROBLEMA CENTRAL",
    Inches(0.6), Inches(0.28), Inches(10), Inches(0.45),
    size=13, bold=True, color=CELESTE)
txt(s, "Dos audiencias completamente\ndistintas. Un solo sitio para ambas.",
    Inches(0.6), Inches(0.72), Inches(10.5), Inches(1.5),
    size=38, bold=True, color=BLANCO)

# Audiencia A
rect(s, Inches(0.5), Inches(2.45), Inches(5.8), Inches(0.48), CELESTE)
txt(s, "AUDIENCIA A — AFILIADO ACTUAL",
    Inches(0.6), Inches(2.49), Inches(5.6), Inches(0.4),
    size=13, bold=True, color=VIOLETA)

aud_a = [
    ("Necesita",  "Sacar turno · Pagar cuota · Número de emergencia · Ver horarios"),
    ("Entra por", "Búsqueda de marca: 'Arevalo turnos' / 'Arevalo pagar cuota'"),
    ("CTA ideal", '"Reservá tu turno"  /  "Acceso afiliados"'),
    ("Frustración actual", "Menú de 6 ítems para encontrar un número de teléfono urgente"),
]
for ii, (k, v) in enumerate(aud_a):
    yy = Inches(3.02) + ii * Inches(0.55)
    bg = VIOLETA_OSC if ii % 2 == 0 else VIOLETA_MED
    rect(s, Inches(0.5), yy, Inches(5.8), Inches(0.51), bg)
    txt(s, k + ":", Inches(0.65), yy + Inches(0.1), Inches(1.4), Inches(0.35),
        size=10, bold=True, color=CELESTE)
    txt(s, v, Inches(2.1), yy + Inches(0.1), Inches(4.1), Inches(0.35),
        size=10, color=BLANCO)

# Audiencia B
rect(s, Inches(7.0), Inches(2.45), Inches(5.8), Inches(0.48), ROSA)
txt(s, "AUDIENCIA B — POTENCIAL NUEVO AFILIADO",
    Inches(7.1), Inches(2.49), Inches(5.6), Inches(0.4),
    size=13, bold=True, color=VIOLETA)

aud_b = [
    ("Necesita",  "Entender qué cubre · Cuánto cuesta · Si hay sucursal cerca"),
    ("Entra por", "Búsqueda genérica: 'prepaga Tucumán' / 'médico sin obra social'"),
    ("CTA ideal", '"Quiero afiliarme"  /  "Consultá precio"'),
    ("Frustración actual", "Mismo menú que el afiliado. Sin precios. Sin diferenciadores."),
]
for ii, (k, v) in enumerate(aud_b):
    yy = Inches(3.02) + ii * Inches(0.55)
    bg = RGBColor(0xCC, 0x44, 0x88) if ii % 2 == 0 else RGBColor(0xBB, 0x33, 0x77)
    rect(s, Inches(7.0), yy, Inches(5.8), Inches(0.51), bg)
    txt(s, k + ":", Inches(7.15), yy + Inches(0.1), Inches(1.4), Inches(0.35),
        size=10, bold=True, color=CELESTE)
    txt(s, v, Inches(8.6), yy + Inches(0.1), Inches(4.1), Inches(0.35),
        size=10, color=BLANCO)

# Solución
rect(s, Inches(0.5), Inches(5.34), W - Inches(1), Inches(0.78), ROSA)
txt(s, "SOLUCIÓN:  La home necesita bifurcar desde el primer segundo.",
    Inches(0.7), Inches(5.4), W - Inches(1.4), Inches(0.32),
    size=14, bold=True, color=VIOLETA)
txt(s, '"¿Ya sos afiliado?"  →  Portal afiliados (turnos, pagos, emergencias)'
       '          |          '
       '"¿Querés conocernos?"  →  Funnel de captación (cobertura, precio, afiliación)',
    Inches(0.7), Inches(5.7), W - Inches(1.4), Inches(0.35),
    size=11, italic=True, color=VIOLETA)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — ARQUITECTURA DE INFORMACIÓN PROPUESTA
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "ARQUITECTURA DE INFORMACIÓN PROPUESTA",
       "De un menú de 15 ítems a un recorrido por audiencia")

# Actual
rect(s, Inches(0.4), Inches(1.18), Inches(5.8), Inches(0.4), ROJO_NOK)
txt(s, "ESTRUCTURA ACTUAL — Confusa para ambas audiencias",
    Inches(0.5), Inches(1.22), Inches(5.6), Inches(0.35),
    size=11, bold=True, color=BLANCO)

actual = [
    "Inicio",
    "Beneficios  ▸  Medicina / Emergencias / Farmacia / Sepelios / Complementos / Servicios",
    "Seguros  ▸  Salud / Renta Diaria / Vida / Accidentes",
    "Afiliación  ▸  Pre-Afiliación / Asesores / Contacto Afiliación",
    "Noticias",
    "Contacto  ▸  Números Útiles / Sucursales",
]
for ii, item in enumerate(actual):
    yy = Inches(1.65) + ii * Inches(0.5)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(0.4), yy, Inches(5.8), Inches(0.46), bg)
    rect(s, Inches(0.4), yy, Inches(0.06), Inches(0.46), ROJO_NOK)
    txt(s, item, Inches(0.56), yy + Inches(0.08), Inches(5.6), Inches(0.32),
        size=10, color=NEGRO)

# Flecha central
txt(s, "→", Inches(6.2), Inches(3.5), Inches(0.8), Inches(0.6),
    size=32, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)

# Propuesta
rect(s, Inches(7.1), Inches(1.18), Inches(5.8), Inches(0.4), VERDE_OK)
txt(s, "ESTRUCTURA PROPUESTA — Bifurcada por audiencia",
    Inches(7.2), Inches(1.22), Inches(5.6), Inches(0.35),
    size=11, bold=True, color=BLANCO)

# Sub header nuevo afiliado
rect(s, Inches(7.1), Inches(1.64), Inches(5.8), Inches(0.3), VIOLETA)
txt(s, "NUEVO AFILIADO", Inches(7.2), Inches(1.66), Inches(5.6), Inches(0.26),
    size=10, bold=True, color=CELESTE)
nuevos = [
    "Cobertura médica  (24 especialidades · sucursales)",
    "Planes y precios  (rangos orientativos + comparativa)",
    "Por qué Arevalo  (diferenciadores + testimonios)",
    "Cómo afiliarse  (3 pasos simples + formulario)",
]
for ii, item in enumerate(nuevos):
    yy = Inches(2.0) + ii * Inches(0.46)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(7.1), yy, Inches(5.8), Inches(0.42), bg)
    rect(s, Inches(7.1), yy, Inches(0.06), Inches(0.42), VIOLETA)
    txt(s, item, Inches(7.26), yy + Inches(0.07), Inches(5.6), Inches(0.3),
        size=10, color=NEGRO)

# Sub header afiliado actual
rect(s, Inches(7.1), Inches(4.0), Inches(5.8), Inches(0.3), ROSA)
txt(s, "AFILIADO ACTUAL", Inches(7.2), Inches(4.02), Inches(5.6), Inches(0.26),
    size=10, bold=True, color=VIOLETA)
afiliados = [
    "Sacá tu turno  (teléfono + app)",
    "Pagá tu cuota  (débito + online)",
    "Emergencias  (número siempre visible)",
    "Sucursales y horarios  (mapa interactivo)",
]
for ii, item in enumerate(afiliados):
    yy = Inches(4.36) + ii * Inches(0.46)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(7.1), yy, Inches(5.8), Inches(0.42), bg)
    rect(s, Inches(7.1), yy, Inches(0.06), Inches(0.42), ROSA)
    txt(s, item, Inches(7.26), yy + Inches(0.07), Inches(5.6), Inches(0.3),
        size=10, color=NEGRO)

divider(s, Inches(6.6))
txt(s,
    "El objetivo: 0 clics para llegar a la acción principal de cada audiencia. Todo lo demás, en segundo nivel.",
    Inches(0.5), Inches(6.68), W - Inches(1), Inches(0.38),
    size=12, bold=True, italic=True, color=VIOLETA, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — ANÁLISIS DE CONTENIDO
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "ANÁLISIS DE CONTENIDO",
       "Qué dice el sitio, qué debería decir y qué falta")

# Propuesta de valor actual vs propuesta
rect(s, Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.38), VIOLETA)
txt(s, "PROPUESTA DE VALOR", Inches(0.5), Inches(1.2), Inches(5), Inches(0.34),
    size=11, bold=True, color=BLANCO)

rect(s, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.9), GRIS_CLARO)
txt(s, "ACTUAL", Inches(0.5), Inches(1.62), Inches(5.8), Inches(0.28),
    size=9, bold=True, color=GRIS_MED)
txt(s, '"En salud tu tiempo cuenta"  —  Genérico. Podría ser de cualquier prepaga del país.',
    Inches(0.5), Inches(1.88), Inches(5.8), Inches(0.55),
    size=11, italic=True, color=NEGRO)

rect(s, Inches(6.5), Inches(1.6), Inches(6.4), Inches(0.9), RGBColor(0xEE, 0xE6, 0xFF))
txt(s, "PROPUESTA", Inches(6.6), Inches(1.62), Inches(6.2), Inches(0.28),
    size=9, bold=True, color=VIOLETA)
txt(s, '"La cobertura médica completa de Tucumán — en tu ciudad, con tu médico, 24 horas."',
    Inches(6.6), Inches(1.88), Inches(6.2), Inches(0.55),
    size=11, bold=True, italic=True, color=VIOLETA)

# Contenido faltante vs existente
rect(s, Inches(0.4), Inches(2.65), Inches(6.0), Inches(0.38), ROJO_NOK)
txt(s, "CONTENIDO FALTANTE CRÍTICO",
    Inches(0.5), Inches(2.68), Inches(5.8), Inches(0.32),
    size=11, bold=True, color=BLANCO)

faltante = [
    ("Precios orientativos",        "La pregunta #1 del potencial afiliado es '¿cuánto sale?'. Sin respuesta = se va a Google a buscarla, y la competencia está ahí."),
    ("Testimonios de afiliados",    "En salud, el boca a boca digital es el driver de conversión más poderoso. Hoy no hay ninguno visible."),
    ("Comparativa de planes",       "Individual / Familiar / Grupal. Una tabla simple que responda '¿qué plan me conviene?' reduce enormemente la fricción."),
    ("Trayectoria con datos",       "'25 años · +X afiliados · 24 especialidades · 5 sucursales'. Esos números generan confianza inmediata."),
    ("Proceso de afiliación visual","El usuario no sabe cuánto tarda ni qué necesita. 3 pasos con íconos resuelven eso."),
]
for ii, (titulo, desc) in enumerate(faltante):
    yy = Inches(3.1) + ii * Inches(0.72)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(0.4), yy, Inches(6.0), Inches(0.68), bg)
    rect(s, Inches(0.4), yy, Inches(0.06), Inches(0.68), ROJO_NOK)
    txt(s, titulo, Inches(0.56), yy + Inches(0.06), Inches(5.7), Inches(0.26),
        size=10, bold=True, color=NEGRO)
    txt(s, desc, Inches(0.56), yy + Inches(0.3), Inches(5.7), Inches(0.34),
        size=9.5, italic=True, color=GRIS_TEXTO)

# Contenido existente que hay que potenciar
rect(s, Inches(6.7), Inches(2.65), Inches(6.2), Inches(0.38), VERDE_OK)
txt(s, "CONTENIDO EXISTENTE A POTENCIAR",
    Inches(6.8), Inches(2.68), Inches(6.0), Inches(0.32),
    size=11, bold=True, color=BLANCO)

potenciar = [
    ("App propia (Google Play)",   "Está disponible pero nunca se promociona. Un banner en el hero + campaña de installs la convierte en diferencial."),
    ("FAQ / Preguntas frecuentes", "Contenido que ya existe pero no está en el menú. Es SEO gratuito si se estructura bien."),
    ("Blog de noticias",           "2 posts en 2026. Con 2 artículos/semana de salud preventiva se puede posicionar en búsquedas locales en 60 días."),
    ("Sucursales interior",        "Concepción, Monteros, Aguilares, Alberdi. Cada una merece su propia página con SEO local."),
    ("Especialidades (24)",        "Una tabla de médicos y horarios no comunica el valor. Hay que convertirlo en 'conocé a tu equipo médico'."),
]
for ii, (titulo, desc) in enumerate(potenciar):
    yy = Inches(3.1) + ii * Inches(0.72)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(6.7), yy, Inches(6.2), Inches(0.68), bg)
    rect(s, Inches(6.7), yy, Inches(0.06), Inches(0.68), VERDE_OK)
    txt(s, titulo, Inches(6.86), yy + Inches(0.06), Inches(6.0), Inches(0.26),
        size=10, bold=True, color=NEGRO)
    txt(s, desc, Inches(6.86), yy + Inches(0.3), Inches(6.0), Inches(0.34),
        size=9.5, italic=True, color=GRIS_TEXTO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — BASES PARA CAMPAÑAS DE ADS Y REDES
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

txt(s, "BASES PARA CAMPAÑAS — ADS Y REDES SOCIALES",
    Inches(0.6), Inches(0.22), Inches(11), Inches(0.48),
    size=14, bold=True, color=CELESTE)
txt(s, "Qué canales, qué audiencias y qué mensaje.",
    Inches(0.6), Inches(0.7), Inches(11), Inches(0.62),
    size=28, bold=True, color=BLANCO)

# Segmentos
rect(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(0.38), ROSA)
txt(s, "SEGMENTOS DE AUDIENCIA PRIORITARIOS",
    Inches(0.5), Inches(1.48), Inches(12.3), Inches(0.32),
    size=11, bold=True, color=VIOLETA)

seg_cols = ["SEGMENTO", "PERFIL", "CANAL IDEAL", "MENSAJE CLAVE"]
seg_cw = [Inches(2.8), Inches(3.5), Inches(2.5), Inches(3.5)]
seg_cx = [Inches(0.4), Inches(3.24), Inches(6.78), Inches(9.32)]
for h_, x, w in zip(seg_cols, seg_cx, seg_cw):
    rect(s, x, Inches(1.9), w - Inches(0.04), Inches(0.36), VIOLETA_OSC)
    txt(s, h_, x + Inches(0.08), Inches(1.92), w - Inches(0.18), Inches(0.3),
        size=9, bold=True, color=CELESTE, align=PP_ALIGN.CENTER)

segs = [
    ("Sin cobertura médica",     "25-55 años, Tucumán capital",   "Google Search + Meta",   '"Cobertura médica completa. Desde $X/mes."'),
    ("Familias con hijos",        "Padres 28-45, salud familiar",  "Meta (FB/IG Stories)",   '"Pediatría, odonto y urgencias para toda la familia."'),
    ("Monotributistas",           "Sin obra social patronal",      "Google Search",          '"Tu prepaga sin obra social. Todo Tucumán cubierto."'),
    ("Interior de la provincia",  "Concepción, Monteros, etc.",    "Meta local por ciudad",  '"Ahora en [ciudad] — la misma atención cerca tuyo."'),
    ("Visitas web sin conversión","Ya visitaron el sitio",         "Meta Retargeting",       '"Todavía podés consultarnos. Es gratis y sin compromiso."'),
]
for ii, row in enumerate(segs):
    yy = Inches(2.32) + ii * Inches(0.52)
    bg = VIOLETA_OSC if ii % 2 == 0 else VIOLETA_MED
    for val, x, w in zip(row, seg_cx, seg_cw):
        rect(s, x, yy, w - Inches(0.04), Inches(0.48), bg)
        txt(s, val, x + Inches(0.08), yy + Inches(0.07),
            w - Inches(0.18), Inches(0.36), size=10, color=BLANCO)

# Pilares de contenido
rect(s, Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.38), CELESTE)
txt(s, "PILARES DE CONTENIDO — INSTAGRAM Y FACEBOOK",
    Inches(0.5), Inches(5.08), Inches(12.3), Inches(0.32),
    size=11, bold=True, color=VIOLETA)

pilares = [
    ("🩺 Salud\npreventiva", "Tips semanales\nde salud"),
    ("👨‍⚕️ Nuestro\nequipo médico", "Presentación de\nespecialistas"),
    ("🏙️ Vida en\nsucursales", "Fotos reales del\nequipo en cada ciudad"),
    ("📲 App +\nturnos", "Cómo usar la app\nen 1 minuto"),
    ("💬 Testimonio\nde afiliado", "Historia real\ncada semana"),
    ("🚨 Emergencias\ny guardia", "Recordatorio del\nnúmero urgente"),
]
px = [Inches(0.4 + i * 2.15) for i in range(6)]
for (titulo, desc), x in zip(pilares, px):
    rect(s, x, Inches(5.5), Inches(2.08), Inches(1.62), RGBColor(0x22, 0x00, 0x99))
    rect(s, x, Inches(5.5), Inches(2.08), Inches(0.06), CELESTE)
    txt(s, titulo, x + Inches(0.12), Inches(5.58), Inches(1.9), Inches(0.55),
        size=11, bold=True, color=CELESTE)
    txt(s, desc, x + Inches(0.12), Inches(6.1), Inches(1.9), Inches(0.7),
        size=10, italic=True, color=BLANCO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — PLAN DE ACCIÓN PRIORIZADO
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "PLAN DE ACCIÓN PRIORIZADO",
       "Ordenado por impacto en negocio · Nivel de esfuerzo")

col_titles = ["#", "HALLAZGO", "PRIORIDAD", "ESFUERZO", "ACCIÓN", "KPI"]
col_x = [Inches(0.38), Inches(0.75), Inches(3.45), Inches(5.05), Inches(6.25), Inches(11.0)]
col_w = [Inches(0.35), Inches(2.65), Inches(1.55), Inches(1.15), Inches(4.7), Inches(2.1)]
for h_, x, w in zip(col_titles, col_x, col_w):
    rect(s, x, Inches(1.15), w, Inches(0.38), VIOLETA)
    txt(s, h_, x + Inches(0.06), Inches(1.17), w - Inches(0.12), Inches(0.32),
        size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

acciones = [
    ("#1", "GA4 apagado",           "🔴 CRÍTICO", "Muy bajo", "Activar enabledGoogle en config WP",                  "Datos desde día 1"),
    ("#2", "Hero carrusel",          "🔴 CRÍTICO", "Medio",   "Reemplazar por hero estático con 2 CTAs primarios",    "+CVR estimado 25%"),
    ("#3", "Bifurcación audiencias", "🔴 CRÍTICO", "Alto",    "Rediseñar home: afiliado vs. nuevo afiliado",           "+CVR +30%"),
    ("#4", "Eventos Meta Pixel",     "🟠 ALTO",    "Bajo",    "Configurar Lead, Contact, ViewContent vía GTM",         "Retargeting activo"),
    ("#5", "Landing captación",      "🟠 ALTO",    "Medio",   "Landing nuevo afiliado + formulario 3 pasos",           "Leads directos"),
    ("#6", "Precios en seguros",     "🟠 ALTO",    "Bajo",    "Agregar rangos orientativos ('desde $X')",              "-rebote 20%"),
    ("#7", "Google Search B2B",      "🟠 ALTO",    "Medio",   "Keywords: prepaga Tucumán + interior provincia",        "CPL objetivo"),
    ("#8", "Retargeting Meta",       "🟠 ALTO",    "Bajo",    "Audiencia visitantes web (pixel ya activo)",            "Recuperar leads"),
    ("#9", "Testimonios",            "🟡 MEDIO",   "Bajo",    "5-8 historias reales de afiliados con foto",            "+confianza CVR"),
    ("#10","SEO local sucursales",   "🟡 MEDIO",   "Medio",   "Páginas individuales Concepción, Monteros, etc.",       "Tráfico local"),
    ("#11","Blog de salud",          "🟡 MEDIO",   "Medio",   "2 artículos/semana con intent local en Tucumán",        "Orgánico 60 días"),
    ("#12","Promover la App",        "🟢 BAJO",    "Bajo",    "Banner hero + campaña App installs Meta",               "+descargas"),
]

for ii, row in enumerate(acciones):
    yy = Inches(1.6) + ii * Inches(0.43)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    for val, x, w in zip(row, col_x, col_w):
        rect(s, x, yy, w, Inches(0.4), bg)
        c = NEGRO
        if "🔴" in val: c = ROJO_NOK
        elif "🟠" in val: c = NARANJA
        elif "🟡" in val: c = RGBColor(0xAA, 0x88, 0x00)
        elif "🟢" in val: c = VERDE_OK
        b = (val.startswith("#"))
        txt(s, val, x + Inches(0.05), yy + Inches(0.06),
            w - Inches(0.1), Inches(0.3),
            size=9.5, bold=b, color=c)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10A — PROPUESTA COMERCIAL: CÓMO TRABAJAMOS
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

rect(s, W - Inches(4.8), 0, Inches(4.8), H * 0.45, ROSA)
rect(s, W - Inches(3.8), H * 0.45, Inches(3.8), H * 0.55, CELESTE)
rect(s, W - Inches(2.4), Inches(1.5), Inches(2.4), Inches(2.5), VIOLETA_OSC)

txt(s, "PROPUESTA COMERCIAL",
    Inches(0.6), Inches(0.28), Inches(8), Inches(0.4),
    size=13, bold=True, color=CELESTE)
txt(s, "Cómo trabajamos\ncon vos.",
    Inches(0.6), Inches(0.72), Inches(8.2), Inches(1.6),
    size=44, bold=True, color=BLANCO)

pillares_com = [
    ("💰", "Todo incluido en el fee",
     "Modificaciones, actualizaciones, estrategia mensual y soporte — sin costos ocultos ni sorpresas."),
    ("📅", "Fee mensual en USD",
     "Precio estable y previsible. Sin contratos de largo plazo. Podés escalar o ajustar mes a mes."),
    ("🧩", "Modelo modular",
     "Activás lo que necesitás hoy. Agregás módulos cuando tu negocio crece o surgen nuevos objetivos."),
    ("📣", "Pauta publicitaria aparte",
     "El presupuesto de Google Ads y Meta Ads es tuyo. Nosotros lo gestionamos — vos decidís cuánto invertir."),
]
positions_com = [
    (Inches(0.5),  Inches(2.65)),
    (Inches(4.75), Inches(2.65)),
    (Inches(0.5),  Inches(4.55)),
    (Inches(4.75), Inches(4.55)),
]
for (icono, titulo, desc), (px, py) in zip(pillares_com, positions_com):
    rect(s, px, py, Inches(3.9), Inches(1.65), VIOLETA_OSC)
    rect(s, px, py, Inches(0.08), Inches(1.65), CELESTE)
    txt(s, icono + "  " + titulo,
        px + Inches(0.2), py + Inches(0.18), Inches(3.55), Inches(0.38),
        size=13, bold=True, color=CELESTE)
    txt(s, desc,
        px + Inches(0.2), py + Inches(0.62), Inches(3.55), Inches(0.9),
        size=11, italic=True, color=BLANCO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10B — PLANES BOREAL
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "PLANES BOREAL",
       "Fee mensual en USD · Todo incluido · Pauta publicitaria aparte")

planes_data = [
    {
        "nombre": "STARTER",
        "precio": "$600 USD/mes",
        "accent": CELESTE,
        "text_col": VIOLETA,
        "precio_bg": RGBColor(0x0A, 0x9A, 0xDD),
        "x": Inches(0.38),
        "recomendado": False,
        "features": [
            ("✓", "Gestión de contenido (12 posts/mes)"),
            ("✓", "Diseño de piezas gráficas para el CM"),
            ("✓", "Edición de videos del CM (hasta 4/mes)"),
            ("✓", "Google Maps / Business Profile"),
            ("✓", "Unificación de marca en plataformas"),
            ("✓", "Actualizaciones web (hasta 3hs/mes)"),
            ("✓", "Reporte mensual básico"),
            ("—", "Google Ads + Meta Ads"),
            ("—", "LinkedIn + TikTok activo"),
            ("—", "SEO técnico mensual"),
            ("—", "Email marketing"),
            ("—", "Dashboard en tiempo real"),
        ]
    },
    {
        "nombre": "PRO",
        "precio": "$1.000 USD/mes",
        "accent": ROSA,
        "text_col": BLANCO,
        "precio_bg": RGBColor(0x44, 0x00, 0xFF),
        "x": Inches(4.65),
        "recomendado": True,
        "features": [
            ("✓", "Todo lo de Starter"),
            ("✓", "Google Ads (Search + Display)"),
            ("✓", "Meta Ads (Facebook + Instagram)"),
            ("✓", "LinkedIn + TikTok activo"),
            ("✓", "SEO técnico mensual"),
            ("✓", "Gestión de reputación online"),
            ("✓", "Reporte ejecutivo mensual con KPIs"),
            ("✓", "Actualizaciones web (hasta 6hs/mes)"),
            ("—", "Email marketing"),
            ("—", "Dashboard en tiempo real"),
            ("—", "Consultoría estratégica quincenal"),
            ("—", "Edición pro videos (hasta 8/mes)"),
        ]
    },
    {
        "nombre": "FULL",
        "precio": "$1.600 USD/mes",
        "accent": VIOLETA,
        "text_col": BLANCO,
        "precio_bg": RGBColor(0x22, 0x00, 0x99),
        "x": Inches(8.92),
        "recomendado": False,
        "features": [
            ("✓", "Todo lo de Pro"),
            ("✓", "Edición pro videos CM (hasta 8/mes)"),
            ("✓", "Email marketing + automatizaciones"),
            ("✓", "Dashboard en tiempo real (Looker)"),
            ("✓", "Blog (2 artículos/mes)"),
            ("✓", "Consultoría estratégica quincenal"),
            ("✓", "Actualizaciones web sin límite"),
            ("✓", "Prioridad de soporte"),
            ("✓", "Reunión estratégica mensual"),
            ("✓", "SEO continuo avanzado"),
            ("✓", "Segmentación de audiencias avanzada"),
            ("✓", "Análisis de competencia mensual"),
        ]
    },
]

col_w_plan = Inches(4.0)
for plan in planes_data:
    x = plan["x"]
    ac = plan["accent"]

    # Banner "RECOMENDADO"
    if plan["recomendado"]:
        rect(s, x, Inches(1.06), col_w_plan, Inches(0.22), ROSA)
        txt(s, "★  RECOMENDADO  ★",
            x + Inches(0.05), Inches(1.07), col_w_plan - Inches(0.1), Inches(0.18),
            size=8, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)

    # Header nombre
    rect(s, x, Inches(1.28), col_w_plan, Inches(0.42), ac)
    txt(s, plan["nombre"],
        x + Inches(0.05), Inches(1.3), col_w_plan - Inches(0.1), Inches(0.36),
        size=16, bold=True, color=plan["text_col"], align=PP_ALIGN.CENTER)

    # Precio
    rect(s, x, Inches(1.7), col_w_plan, Inches(0.58), plan["precio_bg"])
    txt(s, plan["precio"],
        x + Inches(0.05), Inches(1.76), col_w_plan - Inches(0.1), Inches(0.44),
        size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    # Features
    for ii, (check, feature) in enumerate(plan["features"]):
        yy = Inches(2.35) + ii * Inches(0.395)
        bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
        rect(s, x, yy, col_w_plan, Inches(0.37), bg)
        rect(s, x, yy, Inches(0.05), Inches(0.37), ac)
        c_check = VERDE_OK if check == "✓" else GRIS_MED
        txt(s, check,
            x + Inches(0.1), yy + Inches(0.06), Inches(0.26), Inches(0.26),
            size=10, bold=(check == "✓"), color=c_check)
        txt(s, feature,
            x + Inches(0.36), yy + Inches(0.06), col_w_plan - Inches(0.44), Inches(0.26),
            size=9.5, color=NEGRO if check == "✓" else GRIS_MED)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10C — MÓDULOS OPCIONALES
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "MÓDULOS OPCIONALES",
       "Se activan cuando el negocio lo pide · Se suman al plan base sin compromisos")

# LEFT — Activaciones y Producción
rect(s, Inches(0.38), Inches(1.18), Inches(6.1), Inches(0.4), ROSA)
txt(s, "🎬  ACTIVACIONES Y PRODUCCIÓN",
    Inches(0.5), Inches(1.21), Inches(5.9), Inches(0.34),
    size=11, bold=True, color=VIOLETA)

mods_izq = [
    ("🌐", "Diseño Web Completo",
     "Rediseño mobile-first bajo manual de marca. \"Menos es más\".",
     "$1.000 con plan\n$1.500 sin plan", "fee único"),
    ("🎬", "Producción Audiovisual Pro",
     "Cámara profesional + drone + edición de estudio. Reels, spots.",
     "Cotizar", "por sesión"),
    ("🎪", "Activaciones en Calle",
     "Diseño de eventos, materiales impresos, roll-ups, punto de venta.",
     "Cotizar", "por evento"),
    ("📸", "Fotografía Profesional",
     "Sesión de equipo médico, sucursales y producto para web y redes.",
     "Cotizar", "por sesión"),
]
for ii, (icono, nombre, desc, precio, tipo) in enumerate(mods_izq):
    yy = Inches(1.66) + ii * Inches(1.2)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(0.38), yy, Inches(6.1), Inches(1.12), bg)
    rect(s, Inches(0.38), yy, Inches(0.08), Inches(1.12), ROSA)
    txt(s, icono + "  " + nombre,
        Inches(0.58), yy + Inches(0.1), Inches(4.2), Inches(0.32),
        size=11, bold=True, color=VIOLETA)
    txt(s, desc,
        Inches(0.58), yy + Inches(0.42), Inches(4.2), Inches(0.42),
        size=10, italic=True, color=GRIS_TEXTO)
    rect(s, Inches(4.8), yy + Inches(0.15), Inches(1.6), Inches(0.56), VIOLETA)
    txt(s, precio,
        Inches(4.82), yy + Inches(0.17), Inches(1.56), Inches(0.32),
        size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txt(s, tipo,
        Inches(4.82), yy + Inches(0.48), Inches(1.56), Inches(0.2),
        size=8, italic=True, color=CELESTE, align=PP_ALIGN.CENTER)

# RIGHT — Tecnología e Integraciones
rect(s, Inches(6.85), Inches(1.18), Inches(6.1), Inches(0.4), CELESTE)
txt(s, "🤖  TECNOLOGÍA E INTEGRACIONES",
    Inches(6.97), Inches(1.21), Inches(5.9), Inches(0.34),
    size=11, bold=True, color=VIOLETA)

mods_der = [
    ("🔍", "Setup Tracking Completo",
     "GA4 + GTM + Meta Pixel + Conversiones Google Ads.",
     "$350 USD", "fee único"),
    ("🔗", "CRM Integration",
     "Setup HubSpot/Zoho, pipeline de leads, automatización de seguimiento.",
     "Cotizar", "fee único + setup"),
    ("🤖", "Chatbot & Atención 24hs",
     "Bot WhatsApp/web, respuestas automáticas, calificación de leads.",
     "Cotizar", "fee mensual"),
    ("🧠", "Agente Comercial AI",
     "Cierre automatizado de consultas, calificación y derivación al equipo.",
     "Cotizar", "fee mensual"),
]
for ii, (icono, nombre, desc, precio, tipo) in enumerate(mods_der):
    yy = Inches(1.66) + ii * Inches(1.2)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(6.85), yy, Inches(6.1), Inches(1.12), bg)
    rect(s, Inches(6.85), yy, Inches(0.08), Inches(1.12), CELESTE)
    txt(s, icono + "  " + nombre,
        Inches(7.05), yy + Inches(0.1), Inches(4.2), Inches(0.32),
        size=11, bold=True, color=VIOLETA)
    txt(s, desc,
        Inches(7.05), yy + Inches(0.42), Inches(4.2), Inches(0.42),
        size=10, italic=True, color=GRIS_TEXTO)
    rect(s, Inches(11.27), yy + Inches(0.15), Inches(1.6), Inches(0.56), VIOLETA)
    txt(s, precio,
        Inches(11.29), yy + Inches(0.17), Inches(1.56), Inches(0.32),
        size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txt(s, tipo,
        Inches(11.29), yy + Inches(0.48), Inches(1.56), Inches(0.2),
        size=8, italic=True, color=CELESTE, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10D — PROPUESTA ESPECÍFICA PARA AREVALO
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

rect(s, W - Inches(4.8), 0, Inches(4.8), H * 0.44, CELESTE)
rect(s, W - Inches(3.8), H * 0.44, Inches(3.8), H * 0.56, ROSA)
rect(s, W - Inches(2.4), Inches(1.5), Inches(2.4), Inches(2.5), VIOLETA_OSC)

txt(s, "NUESTRA PROPUESTA PARA AREVALO",
    Inches(0.6), Inches(0.28), Inches(8.5), Inches(0.4),
    size=13, bold=True, color=CELESTE)
txt(s, "Empezamos con\nlo que más impacta.",
    Inches(0.6), Inches(0.72), Inches(7.8), Inches(1.5),
    size=40, bold=True, color=BLANCO)

# Plan Pro box
rect(s, Inches(0.5), Inches(2.42), Inches(4.0), Inches(3.95), VIOLETA_OSC)
rect(s, Inches(0.5), Inches(2.42), Inches(0.1), Inches(3.95), ROSA)
txt(s, "★  PLAN PRO",
    Inches(0.7), Inches(2.54), Inches(3.7), Inches(0.42),
    size=18, bold=True, color=ROSA)
txt(s, "$1.000 USD/mes",
    Inches(0.7), Inches(3.0), Inches(3.7), Inches(0.55),
    size=26, bold=True, color=BLANCO)
pro_rec = [
    "Gestión de contenido + diseño para redes",
    "Google Ads + Meta Ads gestionados",
    "LinkedIn + TikTok activos y unificados",
    "SEO técnico mensual",
    "Gestión de reputación y reseñas",
    "Reporte ejecutivo mensual con KPIs",
]
for ii, item in enumerate(pro_rec):
    yy = Inches(3.62) + ii * Inches(0.39)
    txt(s, "✓  " + item,
        Inches(0.72), yy, Inches(3.68), Inches(0.34),
        size=10, color=BLANCO)

# Módulo web
rect(s, Inches(4.82), Inches(2.42), Inches(3.8), Inches(1.8), RGBColor(0x44, 0x00, 0xFF))
rect(s, Inches(4.82), Inches(2.42), Inches(0.08), Inches(1.8), CELESTE)
txt(s, "+ DISEÑO WEB COMPLETO",
    Inches(5.02), Inches(2.54), Inches(3.5), Inches(0.34),
    size=12, bold=True, color=CELESTE)
txt(s, "$1.000 USD",
    Inches(5.02), Inches(2.96), Inches(3.5), Inches(0.52),
    size=26, bold=True, color=BLANCO)
txt(s, "Fee único · incluido en el primer mes\nRediseño bajo manual de marca",
    Inches(5.02), Inches(3.54), Inches(3.5), Inches(0.56),
    size=10, italic=True, color=CELESTE)

# Módulo tracking
rect(s, Inches(4.82), Inches(4.38), Inches(3.8), Inches(1.8), RGBColor(0x44, 0x00, 0xFF))
rect(s, Inches(4.82), Inches(4.38), Inches(0.08), Inches(1.8), CELESTE)
txt(s, "+ SETUP TRACKING COMPLETO",
    Inches(5.02), Inches(4.5), Inches(3.5), Inches(0.34),
    size=12, bold=True, color=CELESTE)
txt(s, "$350 USD",
    Inches(5.02), Inches(4.92), Inches(3.5), Inches(0.52),
    size=26, bold=True, color=BLANCO)
txt(s, "Fee único · GA4 + GTM + Meta Pixel\n+ Conversiones Google Ads",
    Inches(5.02), Inches(5.5), Inches(3.5), Inches(0.56),
    size=10, italic=True, color=CELESTE)

# Resumen inversión
rect(s, Inches(8.9), Inches(2.42), Inches(4.0), Inches(3.95), ROSA)
rect(s, Inches(8.9), Inches(2.42), Inches(0.08), Inches(3.95), VIOLETA)
txt(s, "INVERSIÓN INICIAL",
    Inches(9.1), Inches(2.56), Inches(3.7), Inches(0.34),
    size=12, bold=True, color=VIOLETA)
txt(s, "$2.350 USD",
    Inches(9.1), Inches(2.98), Inches(3.7), Inches(0.65),
    size=30, bold=True, color=VIOLETA)
txt(s, "mes 1 (web + tracking + gestión Pro)",
    Inches(9.1), Inches(3.66), Inches(3.7), Inches(0.3),
    size=10, italic=True, color=VIOLETA)

ln2 = s.shapes.add_shape(1, Inches(9.1), Inches(4.06), Inches(3.5), Pt(1.5))
ln2.fill.solid()
ln2.fill.fore_color.rgb = VIOLETA
ln2.line.fill.background()

txt(s, "DESDE MES 2",
    Inches(9.1), Inches(4.18), Inches(3.7), Inches(0.3),
    size=12, bold=True, color=VIOLETA)
txt(s, "$1.000 USD/mes",
    Inches(9.1), Inches(4.54), Inches(3.7), Inches(0.52),
    size=26, bold=True, color=VIOLETA)
txt(s, "Plan Pro · sin costos fijos adicionales",
    Inches(9.1), Inches(5.1), Inches(3.7), Inches(0.3),
    size=10, italic=True, color=VIOLETA)

# Nota pauta
rect(s, Inches(0.5), H - Inches(0.85), W - Inches(1.0), Inches(0.42), VIOLETA_OSC)
txt(s, "⚡  La pauta publicitaria (Google Ads + Meta Ads) es adicional y queda a cargo del cliente. "
       "El presupuesto óptimo se define en la reunión de kickoff.",
    Inches(0.68), H - Inches(0.8), W - Inches(1.4), Inches(0.32),
    size=10, italic=True, color=CELESTE)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — WIREFRAME CONCEPTUAL DE LA NUEVA HOME
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header(s, "CONCEPTO DE REDISEÑO — NUEVA HOME",
       "Estructura propuesta · Mobile-first · Bifurcada por audiencia")

# Mockup simplificado de la home
MX = Inches(0.5)   # mock left
MW = Inches(5.8)   # mock width

# Nav
rect(s, MX, Inches(1.18), MW, Inches(0.32), VIOLETA)
txt(s, "LOGO   |   Cobertura   Sucursales   Afiliarse   🚨 Emergencias",
    MX + Inches(0.1), Inches(1.2), MW - Inches(0.2), Inches(0.28),
    size=8, color=BLANCO)

# Hero
rect(s, MX, Inches(1.52), MW, Inches(1.45), VIOLETA_MED)
txt(s, "HERO: Propuesta de valor + imagen real",
    MX + Inches(0.12), Inches(1.58), MW - Inches(0.25), Inches(0.38),
    size=10, bold=True, color=CELESTE)
txt(s, '"La cobertura médica de Tucumán.\nEn tu ciudad. Con tu médico. 24 horas."',
    MX + Inches(0.12), Inches(1.94), MW - Inches(0.25), Inches(0.45),
    size=9, italic=True, color=BLANCO)
# 2 CTAs
rect(s, MX + Inches(0.12), Inches(2.44), Inches(2.4), Inches(0.34), ROSA)
txt(s, "Quiero afiliarme  →",
    MX + Inches(0.2), Inches(2.49), Inches(2.2), Inches(0.26),
    size=9, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)
rect(s, MX + Inches(2.7), Inches(2.44), Inches(2.4), Inches(0.34), CELESTE)
txt(s, "Ya soy afiliado  →",
    MX + Inches(2.78), Inches(2.49), Inches(2.2), Inches(0.26),
    size=9, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)

# Por qué Arevalo
rect(s, MX, Inches(3.0), MW, Inches(0.28), GRIS_CLARO)
txt(s, "SECCIÓN: Por qué Arevalo — 4 stats (25 años · 24 esp. · 5 sedes · App propia)",
    MX + Inches(0.1), Inches(3.02), MW - Inches(0.2), Inches(0.22),
    size=8, italic=True, color=GRIS_TEXTO)

# Cobertura
rect(s, MX, Inches(3.3), MW, Inches(0.55), GRIS_CLARO)
txt(s, "SECCIÓN: Cobertura médica — Cards por especialidad + botón turno",
    MX + Inches(0.1), Inches(3.38), MW - Inches(0.2), Inches(0.28),
    size=8, italic=True, color=GRIS_TEXTO)

# Sucursales
rect(s, MX, Inches(3.88), MW, Inches(0.42), GRIS_CLARO)
txt(s, "SECCIÓN: Sucursales en toda la provincia — Mapa interactivo",
    MX + Inches(0.1), Inches(3.94), MW - Inches(0.2), Inches(0.25),
    size=8, italic=True, color=GRIS_TEXTO)

# Testimonios
rect(s, MX, Inches(4.33), MW, Inches(0.42), GRIS_CLARO)
txt(s, "SECCIÓN: Testimonios de afiliados — carrusel de historias reales",
    MX + Inches(0.1), Inches(4.39), MW - Inches(0.2), Inches(0.25),
    size=8, italic=True, color=GRIS_TEXTO)

# App
rect(s, MX, Inches(4.78), MW, Inches(0.4), GRIS_CLARO)
txt(s, "SECCIÓN: App Arevalo — Descargala y sacá turnos en 2 minutos",
    MX + Inches(0.1), Inches(4.84), MW - Inches(0.2), Inches(0.25),
    size=8, italic=True, color=GRIS_TEXTO)

# Footer
rect(s, MX, Inches(5.22), MW, Inches(0.4), VIOLETA)
txt(s, "Footer — Emergencias siempre visible · Redes · Datos de contacto",
    MX + Inches(0.1), Inches(5.28), MW - Inches(0.2), Inches(0.25),
    size=8, color=BLANCO)

# Notas de diseño lado derecho
NX = Inches(7.0)
NW = Inches(5.9)

rect(s, NX, Inches(1.18), NW, Inches(0.38), VIOLETA)
txt(s, "NOTAS DE DISEÑO Y ESTRATEGIA",
    NX + Inches(0.1), Inches(1.2), NW - Inches(0.2), Inches(0.32),
    size=11, bold=True, color=BLANCO)

notas = [
    ("🎯", "Hero estático, no carrusel",
     "Una sola imagen de impacto. Un solo mensaje principal. Dos CTAs claros."),
    ("👥", "Bifurcación en el primer scroll",
     "El usuario elige su recorrido desde el hero: nuevo afiliado o afiliado actual."),
    ("📊", "Datos que generan confianza",
     "25 años · X afiliados · 24 especialidades · 5 sucursales. En el hero o inmediatamente abajo."),
    ("🚨", "Emergencias siempre visible",
     "El número 430-7878 debe estar en el header fijo (sticky), visible en cualquier scroll."),
    ("📱", "Mobile-first",
     "El 70%+ del tráfico en salud es desde celular. Diseño desde mobile hacia desktop."),
    ("💬", "Testimonios con foto real",
     "No íconos genéricos. Foto del afiliado + nombre + ciudad + historia breve."),
    ("🗺️", "Mapa de sucursales interactivo",
     "Concepción, Monteros, Aguilares, Alberdi — cada una con dirección y horario."),
]
for ii, (icono, titulo, desc) in enumerate(notas):
    yy = Inches(1.64) + ii * Inches(0.72)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, NX, yy, NW, Inches(0.68), bg)
    rect(s, NX, yy, Inches(0.06), Inches(0.68), CELESTE)
    txt(s, icono + "  " + titulo,
        NX + Inches(0.15), yy + Inches(0.05), NW - Inches(0.3), Inches(0.28),
        size=10, bold=True, color=VIOLETA)
    txt(s, desc,
        NX + Inches(0.15), yy + Inches(0.33), NW - Inches(0.3), Inches(0.3),
        size=9.5, italic=True, color=GRIS_TEXTO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — PRÓXIMOS PASOS Y PROPUESTA BOREAL
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

rect(s, W - Inches(5.5), 0, Inches(5.5), H * 0.5, ROSA)
rect(s, W - Inches(4.5), H * 0.5, Inches(4.5), H * 0.5, CELESTE)
rect(s, W - Inches(3), Inches(1.5), Inches(3), Inches(2.5), VIOLETA_OSC)

txt(s, "EL DIAGNÓSTICO ESTÁ.",
    Inches(0.6), Inches(1.0), Inches(8.5), Inches(0.85),
    size=46, bold=True, color=BLANCO)
txt(s, "El plan también.",
    Inches(0.6), Inches(1.82), Inches(8.5), Inches(0.85),
    size=46, bold=True, color=ROSA)

txt(s, "El próximo paso es una conversación de 30 minutos\npara arrancar con las acciones críticas.",
    Inches(0.6), Inches(2.88), Inches(7.8), Inches(0.82),
    size=15, italic=True, color=BLANCO)

# Qué incluye Boreal
rect(s, Inches(0.5), Inches(3.85), Inches(6.9), Inches(0.38), CELESTE)
txt(s, "LO QUE INCLUYE TRABAJAR CON BOREAL",
    Inches(0.6), Inches(3.88), Inches(6.7), Inches(0.32),
    size=10, bold=True, color=VIOLETA)

servicios = [
    "Rediseño web orientado a conversión",
    "Setup de tracking completo (GA4 + GTM + Pixel)",
    "Campañas Google Ads + Meta Ads con gestión mensual",
    "Estrategia y producción de contenido para redes",
    "Reportes ejecutivos mensuales con KPIs reales",
]
for ii, srv in enumerate(servicios):
    yy = Inches(4.3) + ii * Inches(0.46)
    bg = VIOLETA_OSC if ii % 2 == 0 else VIOLETA_MED
    rect(s, Inches(0.5), yy, Inches(6.9), Inches(0.42), bg)
    rect(s, Inches(0.5), yy, Inches(0.06), Inches(0.42), ROSA)
    txt(s, f"✓  {srv}",
        Inches(0.66), yy + Inches(0.08), Inches(6.6), Inches(0.3),
        size=11, color=BLANCO)

# Contacto
contacto = [
    ("Javier Bertucci",         "CEO — Boreal Marketing"),
    ("+54 9 11 3104-0042",      "WhatsApp"),
    ("jbertucci@boreal.com.ar", "Email"),
]
cy = Inches(4.1)
for val, label in contacto:
    rect(s, Inches(7.7), cy, Inches(5.1), Inches(0.5), VIOLETA_OSC)
    rect(s, Inches(7.7), cy, Inches(0.08), Inches(0.5), CELESTE)
    txt(s, val,
        Inches(7.88), cy + Inches(0.08), Inches(3.0), Inches(0.34),
        size=13, bold=True, color=CELESTE)
    txt(s, label,
        Inches(11.0), cy + Inches(0.12), Inches(1.7), Inches(0.28),
        size=10, italic=True, color=BLANCO, align=PP_ALIGN.RIGHT)
    cy += Inches(0.6)

txt(s, "BOREAL",
    Inches(0.6), H - Inches(0.9), Inches(3), Inches(0.48),
    size=24, bold=True, color=BLANCO)
txt(s, "MARKETING",
    Inches(0.6), H - Inches(0.5), Inches(3), Inches(0.36),
    size=12, bold=True, color=CELESTE)

footer(s)


# ══════════════════════════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════════════════════════
OUTPUT = "/Users/bertuja/claude-boreal/Arevalo_AuditDigital_Boreal_2026.pptx"
prs.save(OUTPUT)
print(f"\n✅  Presentación generada:")
print(f"   {OUTPUT}")
print(f"\n   15 slides — Auditoría digital + Propuesta económica completa")
print(f"   Diseño · Estructura · Contenido · Campañas · Plan de acción")
