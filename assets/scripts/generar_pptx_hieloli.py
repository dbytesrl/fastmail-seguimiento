"""
Generador de presentación PowerPoint — Hieloli × Boreal Marketing
Ejecutar: python3 generar_pptx_hieloli.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PALETA BOREAL ────────────────────────────────────────────────
VIOLETA    = RGBColor(0x33, 0x00, 0xCC)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO      = RGBColor(0x00, 0x00, 0x00)
CELESTE    = RGBColor(0x12, 0xBD, 0xFF)
ROSA       = RGBColor(0xFE, 0x67, 0xB4)
GRIS_CLARO = RGBColor(0xF5, 0xF5, 0xF5)
GRIS_MED   = RGBColor(0xA0, 0xA0, 0xA0)
VIOLETA_OSC = RGBColor(0x22, 0x00, 0x99)
VIOLETA_MED = RGBColor(0x2A, 0x00, 0xBB)

# ─── MEDIDAS ──────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════

def rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h,
        size=18, bold=False, italic=False, color=BLANCO,
        align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Montserrat"
    return tb

def footer(slide, texto="Confidencial · Boreal Marketing × Hieloli · 2026"):
    rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=VIOLETA)
    txt(slide, texto,
        Inches(0.3), H - Inches(0.35), Inches(9), Inches(0.3),
        size=8, color=BLANCO)
    txt(slide, "boreal.com.ar",
        W - Inches(2), H - Inches(0.35), Inches(1.8), Inches(0.3),
        size=8, color=CELESTE, align=PP_ALIGN.RIGHT)

def header_bar(slide, titulo, subtitulo=None):
    rect(slide, 0, 0, W, Inches(1.05), fill=VIOLETA)
    txt(slide, titulo,
        Inches(0.5), Inches(0.18), W - Inches(1), Inches(0.65),
        size=26, bold=True, color=BLANCO)
    if subtitulo:
        txt(slide, subtitulo,
            Inches(0.5), Inches(0.72), W - Inches(1), Inches(0.3),
            size=12, italic=True, color=CELESTE)

def divider(slide, y, color=CELESTE):
    ln = slide.shapes.add_shape(1, Inches(0.5), y, W - Inches(1), Pt(2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

# Bloques geométricos decorativos
rect(s, W - Inches(5.5), 0, Inches(5.5), H * 0.55, ROSA)
rect(s, W - Inches(4.5), H * 0.55, Inches(4.5), H * 0.45, CELESTE)
rect(s, W - Inches(3), Inches(1.5), Inches(3), Inches(2.5), VIOLETA_OSC)

# Agencia
txt(s, "BOREAL MARKETING",
    Inches(0.6), Inches(0.4), Inches(5), Inches(0.45),
    size=11, bold=True, color=CELESTE)
rect(s, Inches(0.6), Inches(0.88), Inches(2.5), Pt(2.5), ROSA)

# Títulos
txt(s, "Lo que construiste\nmerece llegar",
    Inches(0.6), Inches(1.1), Inches(8), Inches(1.8),
    size=46, bold=True, color=BLANCO)
txt(s, "a muchos más.",
    Inches(0.6), Inches(2.85), Inches(8), Inches(0.95),
    size=46, bold=True, color=ROSA)

# Bajada
txt(s,
    "Una mirada estratégica sobre el potencial de crecimiento\ndigital de Hieloli — y cómo activarlo.",
    Inches(0.6), Inches(3.95), Inches(7.8), Inches(0.85),
    size=15, italic=True, color=BLANCO)

# Cliente + fecha
txt(s, "Preparado para HIELOLI · Marzo 2026",
    Inches(0.6), H - Inches(1.05), Inches(7), Inches(0.38),
    size=11, color=GRIS_MED)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — LO QUE HIELOLI CONSTRUYÓ (CELEBRACIÓN)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "UNA EMPRESA CONSTRUIDA PARA DURAR",
           "Antes de hablar de futuro, hay que reconocer lo que ya existe.")

# Hero statement
rect(s, Inches(0.4), Inches(1.2), W - Inches(0.8), Inches(1.0), VIOLETA)
txt(s,
    "Hieloli no es un emprendimiento. Es una operación industrial consolidada con décadas de trayectoria.",
    Inches(0.6), Inches(1.28), W - Inches(1.2), Inches(0.85),
    size=16, bold=True, italic=True, color=CELESTE, align=PP_ALIGN.CENTER)

# 6 logros en 2 columnas
logros = [
    ("🏭", "Dos plantas productivas",  "Pontevedra y Junín — capacidad industrial real."),
    ("🚛", "Flota propia refrigerada", "Distribución directa sin depender de terceros."),
    ("📅", "365 días al año",          "Sin cortes. La confiabilidad que el mercado valora."),
    ("🧊", "Línea completa",           "Hielo, jugos naturales, agua de mesa. Un solo proveedor."),
    ("📋", "Calidad certificada",      "Análisis mensuales. Estándares que pocos en el rubro tienen."),
    ("🤝", "Comodato de freezers",     "Un modelo de negocio que fideliza clientes automáticamente."),
]

lx = [Inches(0.45), Inches(7.0)]
ly = Inches(2.42)
for ii, (icono, titulo, desc) in enumerate(logros):
    x = lx[ii % 2]
    y = ly + (ii // 2) * Inches(1.25)
    bg = GRIS_CLARO if (ii // 2) % 2 == 0 else BLANCO
    rect(s, x, y, Inches(6.1), Inches(1.1), GRIS_CLARO)
    rect(s, x, y, Inches(0.08), Inches(1.1), CELESTE)
    txt(s, icono + "  " + titulo,
        x + Inches(0.2), y + Inches(0.08), Inches(5.7), Inches(0.42),
        size=14, bold=True, color=VIOLETA)
    txt(s, desc,
        x + Inches(0.2), y + Inches(0.52), Inches(5.7), Inches(0.5),
        size=12, color=NEGRO)

divider(s, Inches(6.95), VIOLETA)
txt(s,
    "Esta base es el activo más importante. El marketing digital es la capa que la hace visible para todo el mercado.",
    Inches(0.5), Inches(7.0), W - Inches(1), Inches(0.35),
    size=12, bold=True, italic=True, color=VIOLETA, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — EL POTENCIAL: GENTE QUE BUSCA LO QUE OFRECEN
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "HAY CLIENTES BUSCANDO EXACTAMENTE LO QUE TIENEN",
           "Estas búsquedas ocurren todos los días en Buenos Aires y alrededores.")

txt(s,
    "El producto existe. La demanda existe. Solo hace falta el puente digital que los conecte.",
    Inches(0.5), Inches(1.18), W - Inches(1), Inches(0.38),
    size=13, italic=True, color=GRIS_MED)

# Tabla keywords
cols = ["BÚSQUEDA ACTIVA EN GOOGLE", "TIPO DE CLIENTE", "POTENCIAL", "OPORTUNIDAD"]
cw   = [Inches(4.5), Inches(2.5), Inches(2.2), Inches(3.5)]
cx   = [Inches(0.4), Inches(4.95), Inches(7.5), Inches(9.75)]

for i, (h, x, w) in enumerate(zip(cols, cx, cw)):
    rect(s, x, Inches(1.65), w - Inches(0.05), Inches(0.42), VIOLETA)
    txt(s, h, x + Inches(0.1), Inches(1.67),
        w - Inches(0.2), Inches(0.38),
        size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

rows = [
    ("hielo por mayor Buenos Aires",      "Gastronomía / comercios",   "🔥 MUY ALTO",  "Google Search"),
    ("proveedor hielo para restaurante",  "Restaurantes / bares",      "🔥 MUY ALTO",  "Google Search"),
    ("distribución hielo diaria CABA",    "Locales con volumen",       "🔥 ALTO",       "Google Search"),
    ("hielo para eventos GBA",            "Organizadores de eventos",  "⚡ ALTO",        "Google + Meta"),
    ("comodato freezer hielo comercio",   "Almacenes / kioscos",       "⚡ ALTO",        "Google Search"),
    ("agua bidones por mayor GBA",        "Oficinas / comercios",      "⚡ MEDIO-ALTO",  "Google + Meta"),
    ("hielo costa atlántica temporada",   "Comercios costeros",        "🌊 ESTACIONAL",  "Google + Meta"),
]
for ri, row in enumerate(rows):
    bg = GRIS_CLARO if ri % 2 == 0 else BLANCO
    yr = Inches(2.13) + ri * Inches(0.5)
    for ci, (val, x, w) in enumerate(zip(row, cx, cw)):
        rect(s, x, yr, w - Inches(0.05), Inches(0.46), bg)
        c = VIOLETA if ci == 2 else NEGRO
        b = (ci == 2)
        txt(s, val, x + Inches(0.1), yr + Inches(0.06),
            w - Inches(0.2), Inches(0.36),
            size=11, bold=b, color=c, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(5.8), W - Inches(0.8), Inches(0.75), VIOLETA)
txt(s,
    "Ninguno de estos clientes potenciales llega hoy a Hieloli por búsqueda digital. "
    "Cada uno que no capturan, lo captura otro. La oportunidad es capturarlo ahora que la competencia digital en el rubro es mínima.",
    Inches(0.6), Inches(5.88), W - Inches(1.2), Inches(0.6),
    size=12, color=CELESTE, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — LOS 3 SEGMENTOS CON MAYOR POTENCIAL
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "3 SEGMENTOS. 3 FUENTES DE CRECIMIENTO SOSTENIDO.")

segs = [
    {
        "num": "01",
        "titulo": "GASTRONOMÍA Y COMERCIOS",
        "color_hdr": VIOLETA,
        "color_acc": CELESTE,
        "badge": "Cliente recurrente",
        "puntos": [
            "Restaurantes, bares, hoteles, supermercados",
            "Pedido semanal o diario → alto LTV",
            "El comodato de freezers los fideliza automáticamente",
            "Canal: Google Search + campaña local",
        ],
        "highlight": "1 cliente = 12 meses de facturación continua"
    },
    {
        "num": "02",
        "titulo": "EVENTOS Y PRODUCCIÓN",
        "color_hdr": ROSA,
        "color_acc": VIOLETA,
        "badge": "Alto volumen puntual",
        "puntos": [
            "Productoras, catering, salones de fiestas",
            "Alta urgencia = alta disposición a pagar",
            "Cobertura GBA y Costa = ventaja diferencial",
            "Canal: Google Search + Meta Ads",
        ],
        "highlight": "Un evento puede equivaler a semanas de pedidos"
    },
    {
        "num": "03",
        "titulo": "DISTRIBUIDORES Y MAYORISTAS",
        "color_hdr": CELESTE,
        "color_acc": VIOLETA,
        "badge": "Volumen masivo",
        "puntos": [
            "Distribuidoras de bebidas y alimentos",
            "Pedidos de alto volumen, contrato estable",
            "Agua de mesa y jugos como complemento",
            "Canal: LinkedIn + contacto directo",
        ],
        "highlight": "Un distribuidor puede triplicar el volumen mensual"
    },
]

sx = [Inches(0.4), Inches(4.65), Inches(8.9)]
sw = Inches(4.0)
for seg, x in zip(segs, sx):
    rect(s, x, Inches(1.2), sw, Inches(0.55), seg["color_hdr"])
    num_c = seg["color_acc"] if seg["color_hdr"] != CELESTE else VIOLETA
    txt_c = BLANCO if seg["color_hdr"] != CELESTE else VIOLETA
    txt(s, seg["num"],
        x + Inches(0.12), Inches(1.22), Inches(0.55), Inches(0.5),
        size=28, bold=True, color=num_c)
    txt(s, seg["titulo"],
        x + Inches(0.7), Inches(1.27), sw - Inches(0.8), Inches(0.42),
        size=12, bold=True, color=txt_c)

    # Badge
    rect(s, x, Inches(1.77), sw, Inches(0.3), GRIS_CLARO)
    txt(s, "▸  " + seg["badge"],
        x + Inches(0.15), Inches(1.79), sw - Inches(0.3), Inches(0.26),
        size=10, italic=True, color=VIOLETA)

    for pi, punto in enumerate(seg["puntos"]):
        yp = Inches(2.13) + pi * Inches(0.68)
        bg = GRIS_CLARO if pi % 2 == 0 else BLANCO
        rect(s, x, yp, sw, Inches(0.64), bg)
        txt(s, f"→  {punto}",
            x + Inches(0.12), yp + Inches(0.1), sw - Inches(0.25), Inches(0.5),
            size=11, color=NEGRO)

    rect(s, x, Inches(4.87), sw, Inches(0.65), seg["color_hdr"])
    c2 = BLANCO if seg["color_hdr"] != CELESTE else VIOLETA
    txt(s, "💡  " + seg["highlight"],
        x + Inches(0.12), Inches(4.93), sw - Inches(0.25), Inches(0.52),
        size=11, bold=True, color=c2, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — EL DIFERENCIAL OCULTO (COMODATO)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

txt(s, "EL ACTIVO ESTRATÉGICO MÁS PODEROSO",
    Inches(0.6), Inches(0.28), Inches(11), Inches(0.48),
    size=13, bold=True, color=CELESTE)
txt(s, "Un modelo de negocio\nque la competencia no tiene.",
    Inches(0.6), Inches(0.75), Inches(11), Inches(1.55),
    size=42, bold=True, color=BLANCO)

# Hero comodato
rect(s, Inches(0.5), Inches(2.5), Inches(7.8), Inches(2.1), ROSA)
rect(s, Inches(0.5), Inches(2.5), Inches(0.12), Inches(2.1), CELESTE)
txt(s, "🧊  COMODATO DE FREEZERS — El diferencial que cambia todo",
    Inches(0.72), Inches(2.62), Inches(7.4), Inches(0.5),
    size=18, bold=True, color=VIOLETA)
txt(s,
    "Hieloli instala un freezer en el local del cliente — sin costo.\n"
    "El cliente recibe su hielo siempre del mismo proveedor.\n"
    "Es fidelización automática: el freezer de Hieloli vive dentro del negocio del cliente.",
    Inches(0.72), Inches(3.15), Inches(7.4), Inches(1.25),
    size=14, color=VIOLETA)

txt(s,
    "Este beneficio todavía no está comunicado con la fuerza que merece en digital. "
    "Una vez que se active, se convierte en el CTA más poderoso de toda la campaña.",
    Inches(0.5), Inches(4.78), Inches(7.8), Inches(0.65),
    size=13, italic=True, bold=True, color=CELESTE)

# 3 pilares adicionales
pilares = [
    ("365 días",      "Sin cortes, sin excusas.\nConfiabilidad que los clientes B2B necesitan."),
    ("Calidad cert.", "Análisis mensuales.\nPerfil ideal para cuentas corporativas exigentes."),
    ("Flota propia",  "Control total de la entrega.\nNo dependemos de ningún tercero."),
]
px = [Inches(8.7), Inches(10.5), Inches(12.3)]
for (t, d), x in zip(pilares, px):
    rect(s, x, Inches(2.5), Inches(1.65), Inches(3.2), VIOLETA_OSC)
    rect(s, x, Inches(2.5), Inches(1.65), Inches(0.08), ROSA)
    txt(s, t, x + Inches(0.18), Inches(2.62), Inches(1.4), Inches(0.42),
        size=13, bold=True, color=ROSA)
    txt(s, d, x + Inches(0.18), Inches(3.1), Inches(1.4), Inches(1.4),
        size=11, color=BLANCO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — CÓMO FUNCIONA EL SISTEMA DIGITAL
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "CÓMO EL MARKETING DIGITAL AMPLIFICA LO QUE YA TIENEN",
           "No se trata de cambiar el negocio. Se trata de hacerlo visible.")

# Funnel visual
etapas = [
    ("DESCUBRIMIENTO", "Alguien busca\n'hielo para mi bar\nen CABA'",          VIOLETA,  "Google Ads\nSEO Local\nInstagram"),
    ("CONEXIÓN",       "Llega a la landing\nde Hieloli y ve\nel comodato",      ROSA,     "Landing B2B\nTestimonios\nGarantía calidad"),
    ("CONTACTO",       "Manda WhatsApp\no llena\nel formulario",                CELESTE,  "WhatsApp\nFormulario\nLlamada directa"),
    ("CLIENTE",        "Recibe su freezer\ny empieza a pedir\ncada semana",     VIOLETA_MED, "Relación\ncontinua\n365 días/año"),
]

ew = Inches(2.8)
ex = [Inches(0.45), Inches(3.6), Inches(6.75), Inches(9.9)]
for ii, (et, desc, col, canales) in enumerate(etapas):
    ey = Inches(1.3)
    rect(s, ex[ii], ey, ew, Inches(0.5), col)
    txt_c = VIOLETA if col == CELESTE else BLANCO
    txt(s, et, ex[ii] + Inches(0.1), ey + Inches(0.08),
        ew - Inches(0.2), Inches(0.38),
        size=13, bold=True, color=txt_c, align=PP_ALIGN.CENTER)

    rect(s, ex[ii], ey + Inches(0.5), ew, Inches(1.55), GRIS_CLARO)
    txt(s, desc, ex[ii] + Inches(0.15), ey + Inches(0.58),
        ew - Inches(0.3), Inches(1.4),
        size=12, color=NEGRO, align=PP_ALIGN.CENTER)

    rect(s, ex[ii], ey + Inches(2.1), ew, Inches(1.1), col)
    txt(s, "Vía:", ex[ii] + Inches(0.15), ey + Inches(2.18),
        ew - Inches(0.3), Inches(0.3),
        size=9, bold=True, color=txt_c if col != CELESTE else VIOLETA)
    txt(s, canales, ex[ii] + Inches(0.15), ey + Inches(2.44),
        ew - Inches(0.3), Inches(0.72),
        size=11, color=txt_c if col != CELESTE else VIOLETA, align=PP_ALIGN.CENTER)

    if ii < 3:
        txt(s, "→",
            ex[ii] + ew + Inches(0.1), ey + Inches(0.95),
            Inches(0.45), Inches(0.5),
            size=24, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)

# Resultado final
rect(s, Inches(0.45), Inches(4.62), W - Inches(0.9), Inches(0.7), VIOLETA)
txt(s,
    "El resultado: clientes nuevos llegan solos. El equipo de ventas cierra con leads ya calificados. La operación crece sin cambiar lo que ya funciona.",
    Inches(0.65), Inches(4.72), W - Inches(1.3), Inches(0.52),
    size=13, bold=True, color=CELESTE, align=PP_ALIGN.CENTER)

# Retargeting box
rect(s, Inches(0.45), Inches(5.45), W - Inches(0.9), Inches(0.78), GRIS_CLARO)
rect(s, Inches(0.45), Inches(5.45), Inches(0.1), Inches(0.78), ROSA)
txt(s,
    "BONUS — Retargeting:  los que visitan el sitio y no contactan reciben anuncios recordatorios en Instagram y Google. "
    "Nadie se pierde para siempre.",
    Inches(0.65), Inches(5.55), W - Inches(1.3), Inches(0.6),
    size=12, italic=True, color=VIOLETA)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — PLAN DE 90 DÍAS
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "PLAN DE 90 DÍAS — UN PASO A LA VEZ",
           "Sin interrumpir la operación. Sin riesgos. Con resultados medibles desde el día 1.")

meses = [
    {
        "mes": "MES 1",
        "titulo": "LOS CIMIENTOS",
        "color": VIOLETA,
        "items": [
            "Instalar medición digital (GA4 + Meta Pixel)",
            "Optimizar presencia en Google Maps",
            "Destacar el comodato en la web",
            "Recolectar testimonios de clientes actuales",
            "Crear landing dedicada para gastronomía",
        ],
        "resultado": "Por primera vez, datos reales de cuánta gente busca lo que ofrecen."
    },
    {
        "mes": "MES 2",
        "titulo": "LA ACTIVACIÓN",
        "color": ROSA,
        "items": [
            "Lanzar campaña Google Search B2B",
            "Lanzar campaña Meta Ads (verano / eventos)",
            "Retargeting: recuperar visitas que no contactaron",
            "Publicaciones regulares en Instagram",
            "Primeros reportes semanales de leads",
        ],
        "resultado": "Primeros leads calificados entrando por canales digitales."
    },
    {
        "mes": "MES 3",
        "titulo": "LA ESCALA",
        "color": CELESTE,
        "items": [
            "Optimizar campañas con datos reales",
            "Campaña especial Costa Atlántica",
            "Landing para distribuidores / mayoristas",
            "Automatizar seguimiento de leads por WhatsApp",
            "Reporte ejecutivo con proyección anual",
        ],
        "resultado": "Sistema funcionando. Crecimiento medible mes a mes."
    },
]

mx = [Inches(0.4), Inches(4.65), Inches(8.9)]
mw = Inches(4.0)
for mes, x in zip(meses, mx):
    rect(s, x, Inches(1.22), mw, Inches(0.55), mes["color"])
    c_t = VIOLETA if mes["color"] == CELESTE else BLANCO
    txt(s, f"{mes['mes']} — {mes['titulo']}",
        x + Inches(0.15), Inches(1.27), mw - Inches(0.3), Inches(0.45),
        size=14, bold=True, color=c_t, align=PP_ALIGN.CENTER)

    for ii, item in enumerate(mes["items"]):
        yy = Inches(1.82) + ii * Inches(0.55)
        bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
        rect(s, x, yy, mw, Inches(0.51), bg)
        txt(s, f"✓  {item}",
            x + Inches(0.12), yy + Inches(0.07), mw - Inches(0.25), Inches(0.4),
            size=11, color=NEGRO)

    rect(s, x, Inches(4.62), mw, Inches(0.7), RGBColor(0xEE, 0xE6, 0xFF))
    rect(s, x, Inches(4.62), Inches(0.08), Inches(0.7), mes["color"])
    txt(s, mes["resultado"],
        x + Inches(0.2), Inches(4.7), mw - Inches(0.35), Inches(0.56),
        size=11, bold=True, italic=True, color=VIOLETA)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — INVERSIÓN Y RETORNO
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

txt(s, "INVERSIÓN Y RETORNO ESTIMADO",
    Inches(0.6), Inches(0.25), Inches(10), Inches(0.45),
    size=13, bold=True, color=CELESTE)
txt(s, "¿Cuánto se invierte?\n¿Cuánto puede rendir?",
    Inches(0.6), Inches(0.7), Inches(10), Inches(1.5),
    size=38, bold=True, color=BLANCO)

# Lado izquierdo — Inversión
rect(s, Inches(0.5), Inches(2.42), Inches(5.8), Inches(0.48), ROSA)
txt(s, "ESTRUCTURA DE INVERSIÓN MENSUAL",
    Inches(0.6), Inches(2.46), Inches(5.6), Inches(0.4),
    size=13, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)

inv = [
    ("Setup Tracking Completo",  "GA4 + GTM + Meta Pixel + Conversiones Google Ads",   "$350 USD",           True),
    ("Diseño Web Completo",       "Rediseño mobile-first bajo manual de marca Hieloli",  "$1.000 USD",         False),
    ("Gestión mensual Boreal",    "Fee agencia — según plan elegido",                    "desde $600 USD/mes", False),
    ("Pauta publicitaria",        "Google Ads + Meta Ads — presupuesto del cliente",     "a definir",          False),
]
for ii, (concepto, detalle, precio, bonif) in enumerate(inv):
    yy = Inches(2.98) + ii * Inches(0.62)
    bg = VIOLETA_OSC if ii % 2 == 0 else VIOLETA_MED
    rect(s, Inches(0.5), yy, Inches(5.8), Inches(0.58), bg)
    txt(s, concepto, Inches(0.65), yy + Inches(0.07), Inches(3.2), Inches(0.28),
        size=12, bold=True, color=BLANCO)
    txt(s, detalle, Inches(0.65), yy + Inches(0.33), Inches(3.2), Inches(0.22),
        size=9, italic=True, color=GRIS_MED)
    if bonif:
        txt(s, precio, Inches(4.1), yy + Inches(0.02), Inches(2.1), Inches(0.24),
            size=10, italic=True, color=GRIS_MED, align=PP_ALIGN.RIGHT)
        # línea tachado visual
        ln_t = s.shapes.add_shape(1, Inches(4.1), yy + Inches(0.12), Inches(2.1), Pt(1.2))
        ln_t.fill.solid(); ln_t.fill.fore_color.rgb = GRIS_MED; ln_t.line.fill.background()
        rect(s, Inches(4.08), yy + Inches(0.3), Inches(2.14), Inches(0.23), ROSA)
        txt(s, "✓  BONIFICADO", Inches(4.1), yy + Inches(0.31), Inches(2.1), Inches(0.2),
            size=9, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)
    else:
        txt(s, precio, Inches(4.1), yy + Inches(0.1), Inches(2.1), Inches(0.38),
            size=13, bold=True, color=CELESTE, align=PP_ALIGN.RIGHT)

# Lado derecho — ROI
rect(s, Inches(7.0), Inches(2.42), Inches(5.8), Inches(0.48), CELESTE)
txt(s, "RETORNO PROYECTADO — ESCENARIO CONSERVADOR",
    Inches(7.1), Inches(2.46), Inches(5.6), Inches(0.4),
    size=11, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)

roi = [
    ("Nuevos clientes B2B por mes",    "5 a 10 cuentas gastronomía"),
    ("Ticket mensual promedio",         "$80.000 – $150.000 ARS por cuenta"),
    ("Facturación adicional (mes 3)",   "$400.000 – $1.5M ARS / mes"),
    ("LTV por cliente (12 meses)",      "$960.000 – $1.8M ARS por cliente"),
]
for ii, (concepto, valor) in enumerate(roi):
    yy = Inches(2.98) + ii * Inches(0.62)
    bg = RGBColor(0x0E, 0x99, 0xD4) if ii % 2 == 0 else RGBColor(0x0A, 0x88, 0xBB)
    rect(s, Inches(7.0), yy, Inches(5.8), Inches(0.58), bg)
    txt(s, concepto, Inches(7.15), yy + Inches(0.08), Inches(3.3), Inches(0.42),
        size=12, color=BLANCO)
    txt(s, valor, Inches(10.4), yy + Inches(0.08), Inches(2.3), Inches(0.42),
        size=11, bold=True, color=VIOLETA, align=PP_ALIGN.RIGHT)

# Break even
rect(s, Inches(7.0), Inches(5.46), Inches(5.8), Inches(0.58), ROSA)
txt(s, "BREAK-EVEN ESTIMADO:   MES 2",
    Inches(7.1), Inches(5.56), Inches(5.6), Inches(0.4),
    size=16, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)

txt(s,
    "* Proyección basada en benchmarks del vertical (distribución B2B alimentos, Argentina). "
    "Los valores exactos se definen según presupuesto acordado.",
    Inches(0.5), Inches(6.25), W - Inches(1), Inches(0.45),
    size=9, italic=True, color=GRIS_MED, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — POR QUÉ AHORA ES EL MOMENTO
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "EL MOMENTO ES AHORA — LA VENTANA ESTÁ ABIERTA")

razones = [
    {
        "n": "01",
        "titulo": "La competencia digital en el rubro es casi inexistente",
        "texto": (
            "Ningún fabricante de hielo en Argentina tiene una presencia digital bien construida. "
            "El que llega primero con campañas + landing + contenido, posiciona su marca como referente del sector. "
            "Esta ventaja no va a estar disponible para siempre."
        ),
        "color": VIOLETA,
    },
    {
        "n": "02",
        "titulo": "Cada año que pasa sin presencia digital es demanda que se pierde",
        "texto": (
            "Hoy hay personas y empresas buscando exactamente lo que Hieloli ofrece. "
            "Sin campañas ni visibilidad digital, esa demanda va a otro proveedor, "
            "que probablemente no tiene ni la calidad ni la infraestructura de Hieloli."
        ),
        "color": ROSA,
    },
    {
        "n": "03",
        "titulo": "El LTV del cliente B2B hace que cada nuevo cliente valga mucho",
        "texto": (
            "Un restaurante que empieza a pedir hielo cada semana puede representar "
            "entre $960.000 y $1.8M ARS anuales en facturación. "
            "Conseguir 5 clientes así en el primer trimestre ya justifica toda la inversión del año."
        ),
        "color": CELESTE,
    },
]

ry = Inches(1.28)
for r in razones:
    rect(s, Inches(0.4), ry, Inches(0.55), Inches(1.48), r["color"])
    c2 = VIOLETA if r["color"] == CELESTE else BLANCO
    txt(s, r["n"],
        Inches(0.4), ry + Inches(0.4), Inches(0.55), Inches(0.65),
        size=22, bold=True, color=c2, align=PP_ALIGN.CENTER)
    txt(s, r["titulo"],
        Inches(1.1), ry + Inches(0.1), Inches(11.7), Inches(0.45),
        size=16, bold=True, color=VIOLETA)
    txt(s, r["texto"],
        Inches(1.1), ry + Inches(0.58), Inches(11.7), Inches(0.82),
        size=13, color=NEGRO)
    ry += Inches(1.65)

divider(s, Inches(6.25), VIOLETA)
txt(s,
    '"Décadas construyendo una empresa sólida. El próximo paso es que todo el mercado lo sepa."',
    Inches(0.5), Inches(6.35), W - Inches(1), Inches(0.5),
    size=15, bold=True, italic=True, color=VIOLETA, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — QUÉ INCLUYE TRABAJAR CON BOREAL
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

txt(s, "LO QUE INCLUYE TRABAJAR CON BOREAL",
    Inches(0.6), Inches(0.28), Inches(10), Inches(0.45),
    size=13, bold=True, color=CELESTE)
txt(s, "Un equipo completo,\nsin estructuras internas.",
    Inches(0.6), Inches(0.72), Inches(9), Inches(1.45),
    size=40, bold=True, color=BLANCO)

servicios = [
    ("🔧", "Setup y tracking",      "GA4, GTM, Meta Pixel, WhatsApp Business, schema markup"),
    ("🎯", "Paid Media completo",   "Google Ads + Meta Ads: campañas, audiencias, puja, optimización"),
    ("🖥️", "Web y landings",        "Landings B2B optimizadas para conversión con identidad Hieloli"),
    ("📊", "Reportes mensuales",    "Dashboard ejecutivo: leads, CPL, ROAS, proyecciones"),
    ("📱", "Redes y contenido",     "Calendario editorial Instagram/Facebook + gestión de comunidad"),
    ("🤝", "Estrategia continua",   "Reunión mensual + plan trimestral + acceso directo al equipo"),
]

lx = [Inches(0.5), Inches(6.95)]
ly = Inches(2.42)
for ii, (icono, titulo, desc) in enumerate(servicios):
    x = lx[ii % 2]
    y = ly + (ii // 2) * Inches(1.32)
    rect(s, x, y, Inches(6.1), Inches(1.18), VIOLETA_OSC)
    rect(s, x, y, Inches(0.1), Inches(1.18), ROSA)
    txt(s, icono + "  " + titulo,
        x + Inches(0.22), y + Inches(0.1), Inches(5.7), Inches(0.42),
        size=14, bold=True, color=ROSA)
    txt(s, desc,
        x + Inches(0.22), y + Inches(0.56), Inches(5.7), Inches(0.54),
        size=12, color=BLANCO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11A — TRABAJAMOS CON TU CM, NO EN SU LUGAR
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

rect(s, W - Inches(4.8), 0, Inches(4.8), H * 0.44, ROSA)
rect(s, W - Inches(3.8), H * 0.44, Inches(3.8), H * 0.56, CELESTE)
rect(s, W - Inches(2.4), Inches(1.5), Inches(2.4), Inches(2.5), VIOLETA_OSC)

txt(s, "MODELO DE TRABAJO",
    Inches(0.6), Inches(0.28), Inches(8), Inches(0.4),
    size=13, bold=True, color=CELESTE)
txt(s, "Trabajamos con\ntu CM, no en su lugar.",
    Inches(0.6), Inches(0.72), Inches(8.0), Inches(1.55),
    size=40, bold=True, color=BLANCO)

# Columna BOREAL
rect(s, Inches(0.5), Inches(2.5), Inches(5.6), Inches(0.42), CELESTE)
txt(s, "BOREAL  diseña y estrategiza",
    Inches(0.62), Inches(2.53), Inches(5.38), Inches(0.36),
    size=13, bold=True, color=VIOLETA)
boreal_items = [
    "Estrategia mensual de contenido",
    "Brief detallado por pieza (formato, copy, timing)",
    "Agenda de publicaciones con objetivos claros",
    "Diseño de piezas gráficas y templates",
    "Análisis de métricas y ajuste de rumbo",
    "Reunión mensual de agenda con el CM",
]
for ii, item in enumerate(boreal_items):
    yy = Inches(3.0) + ii * Inches(0.5)
    bg = VIOLETA_OSC if ii % 2 == 0 else VIOLETA_MED
    rect(s, Inches(0.5), yy, Inches(5.6), Inches(0.46), bg)
    rect(s, Inches(0.5), yy, Inches(0.07), Inches(0.46), CELESTE)
    txt(s, "→  " + item,
        Inches(0.68), yy + Inches(0.08), Inches(5.3), Inches(0.32),
        size=11, color=BLANCO)

# Columna CM
rect(s, Inches(6.5), Inches(2.5), Inches(5.6), Inches(0.42), ROSA)
txt(s, "Tu CM  ejecuta y publica",
    Inches(6.62), Inches(2.53), Inches(5.38), Inches(0.36),
    size=13, bold=True, color=VIOLETA)
cm_items = [
    "Sube los posts según el brief",
    "Filma el contenido indicado en la agenda",
    "Interactúa con la comunidad",
    "Responde mensajes directos",
    "Reporta novedades y feedback del campo",
    "Participa de la reunión mensual de agenda",
]
for ii, item in enumerate(cm_items):
    yy = Inches(3.0) + ii * Inches(0.5)
    bg = RGBColor(0xCC, 0x44, 0x88) if ii % 2 == 0 else RGBColor(0xBB, 0x33, 0x77)
    rect(s, Inches(6.5), yy, Inches(5.6), Inches(0.46), bg)
    rect(s, Inches(6.5), yy, Inches(0.07), Inches(0.46), ROSA)
    txt(s, "✓  " + item,
        Inches(6.68), yy + Inches(0.08), Inches(5.3), Inches(0.32),
        size=11, color=BLANCO)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11B — PLANES BOREAL
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "PLANES BOREAL",
           "Fee mensual en USD · Todo incluido · Pauta publicitaria aparte")

planes_h = [
    {
        "nombre": "STARTER",
        "precio": "$600 USD/mes",
        "accent": CELESTE,
        "text_col": VIOLETA,
        "precio_bg": RGBColor(0x0A, 0x9A, 0xDD),
        "x": Inches(0.38),
        "features": [
            ("✓", "Estrategia mensual de contenido"),
            ("✓", "Briefs para el CM (formato + copy + timing)"),
            ("✓", "Diseño de piezas gráficas (12 posts/mes)"),
            ("✓", "Templates adaptados por red social"),
            ("✓", "Google Maps / Business Profile"),
            ("✓", "Unificación de marca en plataformas"),
            ("✓", "Actualizaciones web (hasta 3hs/mes)"),
            ("✓", "Reporte mensual básico"),
            ("—", "Reuniones mensuales con CM"),
            ("—", "Google Ads + Meta Ads"),
            ("—", "LinkedIn B2B activo"),
            ("—", "SEO técnico mensual"),
        ]
    },
    {
        "nombre": "PRO",
        "precio": "$1.000 USD/mes",
        "accent": ROSA,
        "text_col": BLANCO,
        "precio_bg": RGBColor(0x44, 0x00, 0xFF),
        "x": Inches(4.65),
        "features": [
            ("✓", "Todo lo de Starter"),
            ("✓", "Reuniones mensuales de agenda con CM"),
            ("✓", "Google Ads (Search B2B + Display)"),
            ("✓", "Meta Ads (Facebook + Instagram)"),
            ("✓", "LinkedIn B2B activo"),
            ("✓", "SEO técnico mensual"),
            ("✓", "Estrategia estacional (costa, eventos)"),
            ("✓", "Gestión de reputación y reseñas"),
            ("✓", "Reporte ejecutivo mensual con KPIs"),
            ("✓", "Actualizaciones web (hasta 6hs/mes)"),
            ("—", "Email marketing B2B"),
            ("—", "Dashboard en tiempo real"),
        ]
    },
    {
        "nombre": "FULL",
        "precio": "$1.600 USD/mes",
        "accent": VIOLETA,
        "text_col": BLANCO,
        "precio_bg": RGBColor(0x22, 0x00, 0x99),
        "x": Inches(8.92),
        "features": [
            ("✓", "Todo lo de Pro"),
            ("✓", "Email marketing B2B + automatizaciones"),
            ("✓", "Dashboard en tiempo real (Looker)"),
            ("✓", "Blog B2B (2 artículos/mes)"),
            ("✓", "Estrategia distribuidores y mayoristas"),
            ("✓", "Consultoría estratégica quincenal"),
            ("✓", "Actualizaciones web sin límite"),
            ("✓", "Prioridad de soporte"),
            ("✓", "Análisis de competencia mensual"),
            ("✓", "Segmentación de audiencias avanzada"),
            ("✓", "Reunión de dirección mensual"),
            ("✓", "Plan trimestral documentado"),
        ]
    },
]

col_w_h = Inches(4.0)
for plan in planes_h:
    x = plan["x"]
    ac = plan["accent"]

    rect(s, x, Inches(1.18), col_w_h, Inches(0.42), ac)
    txt(s, plan["nombre"],
        x + Inches(0.05), Inches(1.2), col_w_h - Inches(0.1), Inches(0.36),
        size=16, bold=True, color=plan["text_col"], align=PP_ALIGN.CENTER)

    rect(s, x, Inches(1.6), col_w_h, Inches(0.58), plan["precio_bg"])
    txt(s, plan["precio"],
        x + Inches(0.05), Inches(1.66), col_w_h - Inches(0.1), Inches(0.44),
        size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    for ii, (check, feature) in enumerate(plan["features"]):
        yy = Inches(2.25) + ii * Inches(0.395)
        bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
        rect(s, x, yy, col_w_h, Inches(0.37), bg)
        rect(s, x, yy, Inches(0.05), Inches(0.37), ac)
        c_check = RGBColor(0x00, 0xB0, 0x60) if check == "✓" else GRIS_MED
        txt(s, check,
            x + Inches(0.1), yy + Inches(0.06), Inches(0.26), Inches(0.26),
            size=10, bold=(check == "✓"), color=c_check)
        txt(s, feature,
            x + Inches(0.36), yy + Inches(0.06), col_w_h - Inches(0.44), Inches(0.26),
            size=9.5, color=NEGRO if check == "✓" else GRIS_MED)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11C — MÓDULOS OPCIONALES
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLANCO)
header_bar(s, "MÓDULOS OPCIONALES",
           "Se activan cuando el negocio lo pide · Se suman al plan base sin compromisos")

# LEFT — Activaciones y Producción
rect(s, Inches(0.38), Inches(1.18), Inches(6.1), Inches(0.4), ROSA)
txt(s, "🎬  ACTIVACIONES Y PRODUCCIÓN",
    Inches(0.5), Inches(1.21), Inches(5.9), Inches(0.34),
    size=11, bold=True, color=VIOLETA)

mods_h_izq = [
    ("🌐", "Diseño Web Completo",
     "Rediseño mobile-first, B2B lead-gen, bajo manual de marca.",
     "$1.000 con plan\n$1.500 sin plan", "fee único"),
    ("🎬", "Producción Audiovisual Pro",
     "Planta, flota, producto — cámara + drone + edición de estudio.",
     "Cotizar", "por sesión"),
    ("🎪", "Activaciones en Calle / Eventos",
     "Presencia en ferias gastronómicas, stands en eventos B2B.",
     "Cotizar", "por evento"),
    ("📸", "Fotografía Profesional",
     "Producto, instalaciones, flota — material para web y redes.",
     "Cotizar", "por sesión"),
]
for ii, (icono, nombre, desc, precio, tipo) in enumerate(mods_h_izq):
    yy = Inches(1.66) + ii * Inches(1.2)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(0.38), yy, Inches(6.1), Inches(1.12), bg)
    rect(s, Inches(0.38), yy, Inches(0.08), Inches(1.12), ROSA)
    txt(s, icono + "  " + nombre,
        Inches(0.58), yy + Inches(0.1), Inches(4.2), Inches(0.32),
        size=11, bold=True, color=VIOLETA)
    txt(s, desc,
        Inches(0.58), yy + Inches(0.42), Inches(4.2), Inches(0.42),
        size=10, italic=True, color=GRIS_MED)
    rect(s, Inches(4.8), yy + Inches(0.15), Inches(1.6), Inches(0.56), VIOLETA)
    txt(s, precio,
        Inches(4.82), yy + Inches(0.17), Inches(1.56), Inches(0.32),
        size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txt(s, tipo,
        Inches(4.82), yy + Inches(0.48), Inches(1.56), Inches(0.2),
        size=8, italic=True, color=CELESTE, align=PP_ALIGN.CENTER)

# RIGHT — Tecnología e Integraciones
rect(s, Inches(6.85), Inches(1.18), Inches(6.1), Inches(0.4), CELESTE)
txt(s, "🤖  TECNOLOGÍA E INTEGRACIONES",
    Inches(6.97), Inches(1.21), Inches(5.9), Inches(0.34),
    size=11, bold=True, color=VIOLETA)

mods_h_der = [
    ("🔍", "Setup Tracking Completo",
     "GA4 + GTM + Meta Pixel + Conversiones Google Ads.",
     "$350 USD", "BONIFICADO"),
    ("🔗", "CRM Integration",
     "Pipeline de leads gastronomía, seguimiento de comodatos.",
     "Cotizar", "fee único + setup"),
    ("🤖", "Chatbot & Atención 24hs",
     "Pedidos urgentes 24hs — ideal para horarios de madrugada.",
     "Cotizar", "fee mensual"),
    ("🧠", "Agente Comercial AI",
     "Calificación automática de leads B2B, derivación al equipo.",
     "Cotizar", "fee mensual"),
]
for ii, (icono, nombre, desc, precio, tipo) in enumerate(mods_h_der):
    yy = Inches(1.66) + ii * Inches(1.2)
    bg = GRIS_CLARO if ii % 2 == 0 else BLANCO
    rect(s, Inches(6.85), yy, Inches(6.1), Inches(1.12), bg)
    rect(s, Inches(6.85), yy, Inches(0.08), Inches(1.12), CELESTE)
    txt(s, icono + "  " + nombre,
        Inches(7.05), yy + Inches(0.1), Inches(4.2), Inches(0.32),
        size=11, bold=True, color=VIOLETA)
    txt(s, desc,
        Inches(7.05), yy + Inches(0.42), Inches(4.2), Inches(0.42),
        size=10, italic=True, color=GRIS_MED)
    if tipo == "BONIFICADO":
        txt(s, precio, Inches(11.27), yy + Inches(0.05), Inches(1.6), Inches(0.22),
            size=9, italic=True, color=GRIS_MED, align=PP_ALIGN.CENTER)
        ln_b = s.shapes.add_shape(1, Inches(11.27), yy + Inches(0.14), Inches(1.6), Pt(1.2))
        ln_b.fill.solid(); ln_b.fill.fore_color.rgb = GRIS_MED; ln_b.line.fill.background()
        rect(s, Inches(11.27), yy + Inches(0.3), Inches(1.6), Inches(0.22), ROSA)
        txt(s, "✓ BONIFICADO", Inches(11.29), yy + Inches(0.31), Inches(1.56), Inches(0.2),
            size=8, bold=True, color=VIOLETA, align=PP_ALIGN.CENTER)
    else:
        rect(s, Inches(11.27), yy + Inches(0.15), Inches(1.6), Inches(0.56), VIOLETA)
        txt(s, precio,
            Inches(11.29), yy + Inches(0.17), Inches(1.56), Inches(0.32),
            size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        txt(s, tipo,
            Inches(11.29), yy + Inches(0.48), Inches(1.56), Inches(0.2),
            size=8, italic=True, color=CELESTE, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — CIERRE / CTA
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, VIOLETA)

# Bloques decorativos
rect(s, W - Inches(5.8), 0, Inches(5.8), H * 0.5, ROSA)
rect(s, W - Inches(4.8), H * 0.5, Inches(4.8), H * 0.5, CELESTE)
rect(s, W - Inches(3.2), Inches(1.2), Inches(3.2), Inches(3), VIOLETA_OSC)

# Texto principal
txt(s, "Lo construiste.",
    Inches(0.6), Inches(1.0), Inches(8.5), Inches(0.95),
    size=52, bold=True, color=BLANCO)
txt(s, "Ahora hagámoslo crecer.",
    Inches(0.6), Inches(1.9), Inches(8.5), Inches(0.95),
    size=40, bold=True, color=ROSA)

txt(s,
    "El siguiente paso es una conversación de 30 minutos\npara definir el plan y el presupuesto.",
    Inches(0.6), Inches(3.05), Inches(7.8), Inches(0.85),
    size=16, italic=True, color=BLANCO)

# Datos contacto
contacto = [
    ("Javier Bertucci",         "CEO — Boreal Marketing"),
    ("+54 9 11 3104-0042",      "WhatsApp directo"),
    ("jbertucci@boreal.com.ar", "Email"),
    ("boreal.com.ar",           "Web"),
]
cy = Inches(4.1)
for val, label in contacto:
    rect(s, Inches(0.5), cy, Inches(6.8), Inches(0.52), VIOLETA_OSC)
    rect(s, Inches(0.5), cy, Inches(0.08), Inches(0.52), CELESTE)
    txt(s, val,
        Inches(0.7), cy + Inches(0.08), Inches(4.0), Inches(0.36),
        size=14, bold=True, color=CELESTE)
    txt(s, label,
        Inches(4.75), cy + Inches(0.12), Inches(2.4), Inches(0.28),
        size=11, italic=True, color=BLANCO, align=PP_ALIGN.RIGHT)
    cy += Inches(0.62)

# Logo texto
txt(s, "BOREAL",
    Inches(0.6), H - Inches(0.95), Inches(3), Inches(0.52),
    size=26, bold=True, color=BLANCO)
txt(s, "MARKETING",
    Inches(0.6), H - Inches(0.5), Inches(3), Inches(0.38),
    size=13, bold=True, color=CELESTE)

footer(s, "Confidencial — Preparado exclusivamente para Hieloli · Boreal Marketing 2026")


# ══════════════════════════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════════════════════════
OUTPUT = "/Users/bertuja/claude-boreal/Hieloli_ProposalBoreal_2026.pptx"
prs.save(OUTPUT)
print(f"\n✅  Presentación generada:")
print(f"   {OUTPUT}")
print(f"\n   14 slides — Identidad Boreal · Propuesta comercial completa")
